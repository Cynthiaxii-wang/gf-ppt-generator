#!/usr/bin/env python3
"""Parse DOCX structure and PPTX layout files from a template directory."""

from __future__ import annotations

import argparse
import hashlib
import json
import re
from io import BytesIO
from pathlib import Path
from typing import Any, Iterator

from docx import Document
from docx.document import Document as DocumentObject
from docx.table import Table
from docx.text.paragraph import Paragraph
from docx.oxml.ns import qn
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import xlsxwriter
from xlsxwriter.utility import xl_cell_to_rowcol


CHINESE_CHAPTER_NUMERALS = "一二三四五六七八九十百零〇两"
CHART_NAMESPACE = "http://schemas.openxmlformats.org/drawingml/2006/chart"


def chinese_number(value: int) -> str:
    digits = "零一二三四五六七八九"
    if value < 10:
        return digits[value]
    if value < 20:
        return "十" + (digits[value % 10] if value % 10 else "")
    if value < 100:
        return (
            digits[value // 10]
            + "十"
            + (digits[value % 10] if value % 10 else "")
        )
    return str(value)


def _chart_formula_location(formula: str) -> tuple[str, int, int] | None:
    """Return the worksheet and zero-based first cell referenced by a chart formula."""
    normalized = formula.strip()
    if "!" not in normalized:
        return None
    sheet_name, cell_range = normalized.rsplit("!", 1)
    sheet_name = sheet_name.strip().strip("'").replace("''", "'")
    sheet_name = re.sub(r"^\[[^\]]+\]", "", sheet_name)
    first_cell = cell_range.split(":", 1)[0].replace("$", "")
    try:
        row, column = xl_cell_to_rowcol(first_cell)
    except (AttributeError, TypeError, ValueError):
        return None
    return sheet_name[:31] or "Sheet1", row, column


def chart_cache_workbook(chart_xml: bytes) -> bytes:
    """Build an embedded workbook at the exact ranges used by a cached chart.

    Many brokerage Word reports link charts to an analyst's external workbook.
    The DOCX still carries every plotted category, series name, and value in the
    chart cache. Writing those values back to their original worksheet/ranges
    makes the transferred PowerPoint chart self-contained and data-editable
    without recreating its visual definition.
    """
    root = etree.fromstring(chart_xml)
    namespace = {"c": CHART_NAMESPACE}
    output = BytesIO()
    workbook = xlsxwriter.Workbook(output, {"in_memory": True})
    worksheets: dict[str, Any] = {}
    formats: dict[str, Any] = {}
    wrote_value = False

    reference_nodes = root.xpath(
        ".//*[c:f and (c:numCache or c:strCache)]",
        namespaces=namespace,
    )
    for reference in reference_nodes:
        formula_nodes = reference.xpath("./c:f/text()", namespaces=namespace)
        if not formula_nodes:
            continue
        location = _chart_formula_location(formula_nodes[0])
        if location is None:
            continue
        sheet_name, start_row, start_column = location
        worksheet = worksheets.get(sheet_name)
        if worksheet is None:
            safe_sheet_name = re.sub(
                r'[\[\]:*?/\\]',
                 '_',
                sheet_name,
            ).strip()
            safe_sheet_name = safe_sheet_name[:31]
            if not safe_sheet_name:
                safe_sheet_name = "Sheet1"

            worksheet = workbook.add_worksheet(safe_sheet_name)
            worksheets[sheet_name] = worksheet

        cache_nodes = reference.xpath("./c:numCache | ./c:strCache", namespaces=namespace)
        if not cache_nodes:
            continue
        cache = cache_nodes[0]
        is_numeric = etree.QName(cache).localname == "numCache"
        format_nodes = cache.xpath("./c:formatCode/text()", namespaces=namespace)
        number_format = format_nodes[0] if format_nodes else None
        cell_format = None
        if is_numeric and number_format:
            cell_format = formats.get(number_format)
            if cell_format is None:
                cell_format = workbook.add_format({"num_format": number_format})
                formats[number_format] = cell_format

        for point in cache.xpath("./c:pt", namespaces=namespace):
            try:
                offset = int(point.get("idx", "0"))
            except ValueError:
                offset = 0
            values = point.xpath("./c:v/text()", namespaces=namespace)
            if not values:
                continue
            value: Any = values[0]
            if is_numeric:
                try:
                    value = float(value)
                except ValueError:
                    pass
            worksheet.write(
                start_row + offset,
                start_column,
                value,
                cell_format,
            )
            wrote_value = True

    if not worksheets:
        workbook.add_worksheet("Sheet1")
    if not wrote_value:
        worksheets.get("Sheet1", workbook.get_worksheet_by_name("Sheet1")).write(0, 0, "")
    workbook.close()
    return output.getvalue()


def heading_classification(
    text: str,
    declared_level: int | None = None,
    style_name: str | None = None,
) -> tuple[int | None, str]:
    """Classify a paragraph using numbering first and Word style only as support."""
    normalized = re.sub(r"\s+", " ", text).strip()
    style = (style_name or "").strip().lower()
    # A table of contents repeats numbered headings but is not report body.
    if style.startswith("toc") or style.startswith("目录"):
        return None, "fallback"
    level_one_patterns = (
        rf"^[{CHINESE_CHAPTER_NUMERALS}]+[、．.]\s*\S+",
        rf"^第[{CHINESE_CHAPTER_NUMERALS}0-9]+章(?:[：:、．.\s]|$)",
    )
    if any(re.match(pattern, normalized) for pattern in level_one_patterns):
        return 1, "numbering"
    if re.match(
        rf"^[（(][{CHINESE_CHAPTER_NUMERALS}]+[）)]\s*\S+",
        normalized,
    ):
        return 2, "numbering"
    # Arabic enumerations such as "4." and "4、" are local points, never chapters.
    if re.match(r"^[0-9]+(?:[、．]|\.(?![0-9]))\s*\S+", normalized):
        return 3, "numbering"

    # Word heading styles are auxiliary evidence only. Long prose accidentally
    # carrying a heading style remains body text.
    heading_like = len(normalized) <= 100 and "\n" not in text
    if heading_like and declared_level in {1, 2, 3}:
        return declared_level, "style"
    return None, "fallback"


def semantic_heading_level(
    text: str,
    declared_level: int | None = None,
    style_name: str | None = None,
) -> int | None:
    """Compatibility wrapper used by the slide and content planners."""
    return heading_classification(text, declared_level, style_name)[0]


def iter_doc_blocks(parent: DocumentObject) -> Iterator[Paragraph | Table]:
    """Yield top-level paragraphs and tables in their document order."""
    for child in parent.element.body.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, parent)
        elif child.tag == qn("w:tbl"):
            yield Table(child, parent)


def paragraph_level(paragraph: Paragraph) -> int | None:
    """Return a heading level based on outline level or common style names."""
    p_pr = paragraph._p.pPr
    if p_pr is not None:
        outline = p_pr.find(qn("w:outlineLvl"))
        if outline is not None:
            try:
                return int(outline.get(qn("w:val"))) + 1
            except (TypeError, ValueError):
                pass

    style_name = paragraph.style.name if paragraph.style is not None else ""
    patterns = (
        r"^Heading\s*(\d+)$",
        r"^标题\s*(\d+)$",
        r"^(\d+)\s*级标题$",
        r"^.*?([一二三四五六七八九十])级标题$",
    )
    for pattern in patterns:
        match = re.match(pattern, style_name, flags=re.IGNORECASE)
        if match:
            token = match.group(1)
            chinese_levels = {
                "一": 1,
                "二": 2,
                "三": 3,
                "四": 4,
                "五": 5,
                "六": 6,
                "七": 7,
                "八": 8,
                "九": 9,
                "十": 10,
            }
            return chinese_levels.get(token, int(token) if token.isdigit() else None)
    return None


MARKER_PATTERN = re.compile(r"^\[Table_[A-Za-z0-9_]+\]$")


def _normalized_marker_text(text: str) -> str:
    return re.sub(r"\s+", "", text)


def marker_value_paragraphs(
    document: DocumentObject,
    marker: str,
    *,
    count: int | None = None,
    stop_prefix: str | None = None,
) -> list[Paragraph]:
    """Read paragraphs following a front-page marker, independent of run splits.

    Brokerage Word files may store ``[Table_Title]`` as one ``w:t`` node or as
    three separate runs (``[``, ``Table_Title``, ``]``). Matching the complete
    paragraph text makes both OOXML representations equivalent.
    """
    paragraph_elements = document.element.xpath(".//w:p")
    target = _normalized_marker_text(marker)
    marker_index: int | None = None
    for index, paragraph_element in enumerate(paragraph_elements):
        text = "".join(paragraph_element.xpath(".//w:t/text()"))
        if _normalized_marker_text(text) == target:
            marker_index = index
            break
    if marker_index is None:
        return []

    values: list[Paragraph] = []
    for paragraph_element in paragraph_elements[marker_index + 1 :]:
        paragraph = Paragraph(paragraph_element, document)
        text = paragraph.text.strip()
        if not text:
            continue
        if MARKER_PATTERN.fullmatch(_normalized_marker_text(text)):
            break
        if stop_prefix and text.startswith(stop_prefix):
            break
        values.append(paragraph)
        if count is not None and len(values) >= count:
            break
    return values


def marker_rows(document: DocumentObject, marker: str, count: int = 1) -> list[list[str]]:
    """Compatibility wrapper returning marker values as text-token rows."""
    return [
        [run.text for run in paragraph.runs if run.text] or [paragraph.text]
        for paragraph in marker_value_paragraphs(document, marker, count=count)
    ]


def _run_uses_emphasis_font(run: Any) -> bool:
    """Recognize Word emphasis stored as a font weight instead of w:b."""
    run_properties = run._r.rPr
    if run_properties is None or run_properties.rFonts is None:
        return False
    font_names = [
        run_properties.rFonts.get(qn(f"w:{attribute}")) or ""
        for attribute in ("ascii", "hAnsi", "eastAsia", "cs")
    ]
    return any(
        re.search(r"(?:bold|semibold|demibold|medium|heavy|black)", name, re.I)
        for name in font_names
    )


def run_is_emphasized(run: Any) -> bool:
    return run.bold is True or _run_uses_emphasis_font(run)


def marker_paragraph_metadata(
    document: DocumentObject,
    marker: str,
    *,
    stop_prefix: str | None = None,
) -> list[dict[str, Any]]:
    """Return rich paragraphs following a marker in the same Word table cell."""
    result: list[dict[str, Any]] = []
    for paragraph in marker_value_paragraphs(
        document,
        marker,
        stop_prefix=stop_prefix,
    ):
        text = paragraph.text.strip()
        runs, _ = paragraph_run_metadata(paragraph)
        result.append({"text": text, "runs": runs})
    return result


def report_metadata(document: DocumentObject) -> dict[str, Any]:
    title_paragraphs = marker_value_paragraphs(document, "[Table_Title]")
    industry_paragraphs = marker_value_paragraphs(
        document,
        "[Table_IndustryAndDate]",
        count=1,
    )
    summary_value_paragraphs = marker_value_paragraphs(
        document,
        "[Table_Summary]",
        stop_prefix="风险提示",
    )
    author_paragraphs = marker_value_paragraphs(document, "[Table_Author]")
    legal_paragraphs = marker_value_paragraphs(document, "[Table_LegalDisclaimer]")
    notice_paragraphs = marker_value_paragraphs(document, "[Table_ImportantNotices]")
    title = "".join(paragraph.text.strip() for paragraph in title_paragraphs)
    industry_and_date = " ".join(
        paragraph.text.strip() for paragraph in industry_paragraphs
    )
    summary = "".join(
        paragraph.text.strip() for paragraph in summary_value_paragraphs
    )
    summary_paragraphs = marker_paragraph_metadata(
        document,
        "[Table_Summary]",
        stop_prefix="风险提示",
    )
    legal_disclaimer = "\n".join(
        paragraph.text.strip() for paragraph in legal_paragraphs
    )
    important_notices = "\n".join(
        paragraph.text.strip() for paragraph in notice_paragraphs
    )
    date_match = re.search(r"\d{4}[/-]\d{1,2}[/-]\d{1,2}", industry_and_date)
    authors: list[str] = []
    for paragraph in author_paragraphs:
        name = paragraph.text.strip()
        if re.fullmatch(r"[\u4e00-\u9fff]{2,4}", name) and name not in authors:
            authors.append(name)
    return {
        "title": title or None,
        "report_date": date_match.group(0) if date_match else None,
        "report_type": industry_and_date or None,
        "summary": summary or None,
        "summary_paragraphs": summary_paragraphs,
        "legal_disclaimer": (
            "\n".join(
                part
                for part in (legal_disclaimer, important_notices)
                if part
            )
            or None
        ),
        "authors": authors,
        "source": "document_marker_table" if title or summary else "fallback",
    }


def document_title(
    document: DocumentObject,
    paragraphs: list[Paragraph],
    metadata: dict[str, Any] | None = None,
    fallback_name: str | None = None,
) -> str:
    """Find the document title, preferring explicit Title-style paragraphs."""
    if metadata and metadata.get("title"):
        return str(metadata["title"])
    for paragraph in paragraphs:
        if paragraph.text.strip() and paragraph.style is not None:
            if paragraph.style.name.lower() in {"title", "标题"}:
                return paragraph.text.strip()

    core_title = (document.core_properties.title or "").strip()
    generic_titles = {"广发报告", "证券研究报告", "研究报告", "报告"}
    if core_title and core_title not in generic_titles:
        return core_title

    if fallback_name:
        fallback = re.sub(r"^【[^】]+】", "", fallback_name).strip()
        fallback = re.sub(r"(?:V\\d+(?:\\.\\d+)?|\\(\\d+\\))$", "", fallback).strip()
        if fallback:
            return fallback

    for paragraph in paragraphs:
        if paragraph.text.strip():
            return paragraph.text.strip()
    return ""


def table_to_dict(table: Table, index: int) -> dict[str, Any]:
    wrapper_rows = [[cell.text for cell in row.cells] for row in table.rows]

    # Brokerage reports often use a one-column wrapper table for the caption
    # and source, with the actual editable data table nested inside its middle
    # cell. python-docx's outer ``cell.text`` intentionally omits nested-table
    # cells, which previously left only caption/source rows and produced an
    # empty PPT slide. Prefer the largest multi-column nested table as the
    # renderable payload while retaining wrapper metadata.
    nested_candidates: list[list[list[str]]] = []
    for nested_element in table._element.xpath(".//w:tbl"):
        nested_table = Table(nested_element, table._parent)
        nested_rows = [
            [cell.text for cell in row.cells]
            for row in nested_table.rows
        ]
        nested_column_count = max(
            (len(row) for row in nested_rows),
            default=0,
        )
        if len(nested_rows) >= 2 and nested_column_count >= 2:
            nested_candidates.append(nested_rows)

    rows = (
        max(
            nested_candidates,
            key=lambda candidate: (
                len(candidate)
                * max((len(row) for row in candidate), default=0)
            ),
        )
        if nested_candidates
        else wrapper_rows
    )
    wrapper_texts = [
        cell.strip()
        for row in wrapper_rows
        for cell in row
        if cell.strip()
    ]
    result = {
        "index": index,
        "row_count": len(rows),
        "column_count": max((len(row) for row in rows), default=0),
        "rows": rows,
    }
    if nested_candidates:
        result["wrapper_rows"] = wrapper_rows
        caption = next(
            (
                text
                for text in wrapper_texts
                if re.match(r"^(?:图|表)\s*\d+\s*[：:]", text)
            ),
            "",
        )
        source_text = next(
            (
                text
                for text in wrapper_texts
                if "数据来源" in text or "资料来源" in text
            ),
            "",
        )
        if caption:
            result["caption"] = caption
        if source_text:
            result["source_text"] = source_text
    return result


def paragraph_run_metadata(paragraph: Paragraph) -> tuple[list[dict[str, Any]], list[str]]:
    runs: list[dict[str, Any]] = []
    bold_buffer: list[str] = []
    bold_segments: list[str] = []
    for run in paragraph.runs:
        text = run.text
        if not text:
            continue
        is_bold = run_is_emphasized(run)
        runs.append(
            {
                "text": text,
                "bold": is_bold,
                "italic": run.italic is True,
            }
        )
        if is_bold:
            bold_buffer.append(text)
        elif bold_buffer:
            segment = "".join(bold_buffer).strip()
            if segment:
                bold_segments.append(segment)
            bold_buffer = []
    if bold_buffer:
        segment = "".join(bold_buffer).strip()
        if segment:
            bold_segments.append(segment)

    bold_sentences: list[str] = []
    for segment in bold_segments:
        bold_sentences.extend(
            sentence.strip()
            for sentence in re.split(r"(?<=[。！？；])", segment)
            if sentence.strip()
        )
    return runs, bold_sentences


def extract_docx_images(document: DocumentObject, output_dir: Path) -> list[dict[str, Any]]:
    image_dir = output_dir / "docx_images"
    image_dir.mkdir(parents=True, exist_ok=True)
    images: list[dict[str, Any]] = []

    # Query DrawingML image references directly. This covers both inline and
    # floating pictures, and safely ignores inline charts/SmartArt objects.
    for index, blip in enumerate(document.element.body.xpath(".//a:blip"), start=1):
        relationship_id = blip.get(qn("r:embed")) or blip.get(qn("r:link"))
        if not relationship_id or relationship_id not in document.part.related_parts:
            continue
        image_part = document.part.related_parts[relationship_id]
        if not image_part.content_type.startswith("image/"):
            continue
        source_name = Path(image_part.partname).name
        digest = hashlib.sha256(image_part.blob).hexdigest()
        target_name = f"{index:03d}_{digest[:12]}_{source_name}"
        target = image_dir / target_name
        if not target.exists():
            target.write_bytes(image_part.blob)
        drawing = blip
        while drawing is not None and drawing.tag not in {
            qn("wp:inline"),
            qn("wp:anchor"),
        }:
            drawing = drawing.getparent()
        extent = drawing.find(qn("wp:extent")) if drawing is not None else None
        images.append(
            {
                "index": index,
                "relationship_id": relationship_id,
                "source_name": source_name,
                "content_type": image_part.content_type,
                "placement": (
                    "inline"
                    if drawing is not None and drawing.tag == qn("wp:inline")
                    else "floating"
                    if drawing is not None
                    else "unknown"
                ),
                "width_emu": int(extent.get("cx")) if extent is not None else None,
                "height_emu": int(extent.get("cy")) if extent is not None else None,
                "extracted_path": target.relative_to(output_dir.parent).as_posix(),
                "sha256": digest,
            }
        )
    return images


def extract_docx_charts(
    document: DocumentObject,
    output_dir: Path,
) -> list[dict[str, Any]]:
    chart_dir = output_dir / "docx_charts"
    chart_dir.mkdir(parents=True, exist_ok=True)
    charts: list[dict[str, Any]] = []
    for index, chart_element in enumerate(
        document.element.xpath(".//c:chart"),
        start=1,
    ):
        relationship_id = chart_element.get(qn("r:id"))
        if not relationship_id:
            continue
        chart_part = document.part.related_parts.get(relationship_id)
        if chart_part is None:
            continue
        drawing = chart_element
        while drawing is not None and drawing.tag not in {
            qn("wp:inline"),
            qn("wp:anchor"),
        }:
            drawing = drawing.getparent()
        extent = drawing.find(qn("wp:extent")) if drawing is not None else None
        target = chart_dir / f"{index:03d}_chart.xml"
        target.write_bytes(chart_part.blob)
        assets: list[dict[str, Any]] = []
        workbook_path: Path | None = None
        workbook_source = "cache"
        external_workbook: str | None = None
        for chart_relationship_id, relationship in chart_part.rels.items():
            if relationship.is_external:
                if relationship.reltype.endswith(("/oleObject", "/package")):
                    external_workbook = relationship.target_ref
                continue
            related_part = relationship.target_part
            suffix = Path(str(related_part.partname)).suffix or ".bin"
            if relationship.reltype.endswith("/package") and suffix.lower() == ".xlsx":
                workbook_path = chart_dir / f"{index:03d}_data.xlsx"
                workbook_path.write_bytes(related_part.blob)
                workbook_source = "embedded"
                continue
            asset_kind = (
                "style"
                if relationship.reltype.endswith("/chartStyle")
                else "colors"
                if relationship.reltype.endswith("/chartColorStyle")
                else "theme"
                if relationship.reltype.endswith("/themeOverride")
                else "related"
            )
            asset_path = chart_dir / f"{index:03d}_{asset_kind}{suffix}"
            asset_path.write_bytes(related_part.blob)
            assets.append(
                {
                    "kind": asset_kind,
                    # Keep the source rId so generate_ppt.py can rewrite every
                    # r:id in the copied chart XML to the relationship ID
                    # allocated in the destination PPTX package.  Relationship
                    # iteration order is not a stable substitute for this map.
                    "relationship_id": chart_relationship_id,
                    "relationship_type": relationship.reltype,
                    "content_type": related_part.content_type,
                    "extracted_path": asset_path.relative_to(output_dir.parent).as_posix(),
                }
            )
        if workbook_path is None:
            workbook_path = chart_dir / f"{index:03d}_data.xlsx"
            workbook_path.write_bytes(chart_cache_workbook(chart_part.blob))
        charts.append(
            {
                "index": index,
                "relationship_id": relationship_id,
                "extracted_path": target.relative_to(output_dir.parent).as_posix(),
                "object_type": "native_chart",
                "transfer_mode": "native_editable",
                "placement": (
                    "inline"
                    if drawing is not None and drawing.tag == qn("wp:inline")
                    else "floating"
                    if drawing is not None
                    else "unknown"
                ),
                "width_emu": int(extent.get("cx")) if extent is not None else None,
                "height_emu": int(extent.get("cy")) if extent is not None else None,
                "workbook_path": workbook_path.relative_to(output_dir.parent).as_posix(),
                "workbook_source": workbook_source,
                "external_workbook": external_workbook,
                "related_assets": assets,
                "sha256": hashlib.sha256(chart_part.blob).hexdigest(),
            }
        )
    return charts


def parse_docx(path: Path, output_dir: Path) -> dict[str, Any]:
    document = Document(path)
    paragraphs = list(document.paragraphs)
    metadata = report_metadata(document)
    title = document_title(document, paragraphs, metadata, path.stem)
    headings_1: list[dict[str, Any]] = []
    headings_2: list[dict[str, Any]] = []
    headings_3: list[dict[str, Any]] = []
    body: list[dict[str, Any]] = []
    toc: list[dict[str, Any]] = []
    tables: list[dict[str, Any]] = []
    content: list[dict[str, Any]] = []
    images = extract_docx_images(document, output_dir)
    charts = extract_docx_charts(document, output_dir)
    paragraph_index = 0
    table_index = 0
    image_cursor = 0
    chart_cursor = 0
    order_index = 0
    heading_1_counter = 0

    def append_content(item: dict[str, Any]) -> None:
        nonlocal order_index
        order_index += 1
        item["order_index"] = order_index
        content.append(item)

    def append_block_images(
        count: int,
        parent_type: str,
        parent_index: int,
    ) -> list[int]:
        nonlocal image_cursor
        indexes: list[int] = []
        for image in images[image_cursor : image_cursor + count]:
            indexes.append(image["index"])
            image_entry = {
                "type": "image",
                **image,
                "parent_type": parent_type,
                "parent_index": parent_index,
            }
            append_content(image_entry)
            image["content_order_index"] = image_entry["order_index"]
            image["parent_type"] = parent_type
            image["parent_index"] = parent_index
        image_cursor += count
        return indexes

    for block in iter_doc_blocks(document):
        if isinstance(block, Paragraph):
            paragraph_index += 1
            text = block.text.strip()
            block_image_count = len(block._p.xpath(".//a:blip"))
            block_chart_count = len(block._p.xpath(".//c:chart"))
            if not text and not block_image_count and not block_chart_count:
                continue
            runs, bold_sentences = paragraph_run_metadata(block)
            if not text:
                append_block_images(
                    block_image_count,
                    "paragraph",
                    paragraph_index,
                )
                if block_chart_count:
                    append_content(
                        {
                            "type": "chart",
                            "index": paragraph_index,
                            "text": "",
                            "heading_level": None,
                            "heading_source": "fallback",
                            "chart_indexes": [
                                chart["index"]
                                for chart in charts[
                                    chart_cursor : chart_cursor + block_chart_count
                                ]
                            ],
                        }
                    )
                chart_cursor += block_chart_count
                continue
            style_name = block.style.name if block.style is not None else None
            level, heading_source = heading_classification(
                text,
                paragraph_level(block),
                style_name,
            )
            heading_text = text
            if level == 1:
                heading_1_counter += 1
                if (
                    heading_source == "style"
                    and semantic_heading_level(heading_text) is None
                ):
                    heading_text = (
                        f"{chinese_number(heading_1_counter)}、{heading_text}"
                    )
            trailing_body = ""
            if level == 1 and heading_source == "numbering" and "\n" in text:
                heading_text, trailing_body = (
                    part.strip() for part in text.split("\n", 1)
                )
            entry = {
                "index": paragraph_index,
                "text": heading_text,
                "style": style_name,
                "heading_level": level,
                "heading_source": heading_source,
                "runs": runs,
                "bold_sentences": bold_sentences,
            }
            if block_image_count:
                entry["image_indexes"] = [
                    image["index"]
                    for image in images[image_cursor : image_cursor + block_image_count]
                ]
            if block_chart_count:
                entry["chart_indexes"] = [
                    chart["index"]
                    for chart in charts[chart_cursor : chart_cursor + block_chart_count]
                ]
            if text == title:
                kind = "title"
            elif (style_name or "").strip().lower().startswith(("toc", "目录")):
                kind = "toc"
                toc.append(entry)
            elif level == 1:
                kind = "heading_1"
                headings_1.append(entry)
            elif level == 2:
                kind = "heading_2"
                headings_2.append(entry)
            elif level == 3:
                kind = "heading_3"
                headings_3.append(entry)
            else:
                kind = "body"
                body.append(entry)
            append_content({"type": kind, **entry})
            if trailing_body:
                trailing_entry = {
                    "index": paragraph_index,
                    "text": trailing_body,
                    "style": style_name,
                    "heading_level": None,
                    "heading_source": "fallback",
                    "runs": [],
                    "bold_sentences": [],
                }
                body.append(trailing_entry)
                append_content({"type": "body", **trailing_entry})
            if block_image_count:
                append_block_images(
                    block_image_count,
                    "paragraph",
                    paragraph_index,
                )
            chart_cursor += block_chart_count
        else:
            table_index += 1
            table_entry = table_to_dict(block, table_index)
            block_image_count = len(block._element.xpath(".//a:blip"))
            block_chart_count = len(block._element.xpath(".//c:chart"))
            if block_image_count:
                table_entry["image_indexes"] = [
                    image["index"]
                    for image in images[image_cursor : image_cursor + block_image_count]
                ]
            if block_chart_count:
                table_entry["chart_indexes"] = [
                    chart["index"]
                    for chart in charts[chart_cursor : chart_cursor + block_chart_count]
                ]
            tables.append(table_entry)
            append_content({"type": "table", **table_entry})
            if block_image_count:
                append_block_images(
                    block_image_count,
                    "table",
                    table_index,
                )
            chart_cursor += block_chart_count

    return {
        "file": path.name,
        "title": title,
        "report_metadata": metadata,
        "heading_1": headings_1,
        "heading_2": headings_2,
        "heading_3": headings_3,
        "body": body,
        "toc": toc,
        "tables": tables,
        "images": images,
        "charts": charts,
        "content_order": content,
    }


def shape_type_name(shape: Any) -> str:
    try:
        return MSO_SHAPE_TYPE(shape.shape_type).name
    except (TypeError, ValueError):
        return str(shape.shape_type)


def shape_geometry(shape: Any) -> dict[str, Any]:
    return {
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
        "unit": "EMU",
    }


def shape_text(shape: Any) -> str | None:
    if getattr(shape, "has_text_frame", False):
        return shape.text
    if getattr(shape, "has_table", False):
        return "\n".join(
            "\t".join(cell.text for cell in row.cells) for row in shape.table.rows
        )
    return None


def serialize_shape(shape: Any, index: int, parent_group: str | None = None) -> dict[str, Any]:
    item: dict[str, Any] = {
        "index": index,
        "name": shape.name,
        "shape_type": shape_type_name(shape),
        "shape_type_code": int(shape.shape_type),
        **shape_geometry(shape),
        "text": shape_text(shape),
        "is_picture": shape.shape_type in {
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.LINKED_PICTURE,
        },
    }
    if parent_group is not None:
        item["parent_group"] = parent_group
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        item["children"] = [
            serialize_shape(child, child_index, shape.name)
            for child_index, child in enumerate(shape.shapes, start=1)
        ]
    return item


def flatten_shapes(items: list[dict[str, Any]]) -> Iterator[dict[str, Any]]:
    for item in items:
        yield item
        yield from flatten_shapes(item.get("children", []))


def parse_pptx(path: Path) -> dict[str, Any]:
    presentation = Presentation(path)
    slides: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        shapes = [
            serialize_shape(shape, index)
            for index, shape in enumerate(slide.shapes, start=1)
        ]
        pictures = [
            {
                "name": shape["name"],
                "left": shape["left"],
                "top": shape["top"],
                "width": shape["width"],
                "height": shape["height"],
                "unit": "EMU",
            }
            for shape in flatten_shapes(shapes)
            if shape["is_picture"]
        ]
        slides.append(
            {
                "slide_number": slide_number,
                "shape_count": len(list(flatten_shapes(shapes))),
                "shapes": shapes,
                "picture_frames": pictures,
            }
        )
    return {
        "file": path.name,
        "slide_count": len(presentation.slides),
        "slide_width": presentation.slide_width,
        "slide_height": presentation.slide_height,
        "unit": "EMU",
        "slides": slides,
    }


def write_json(path: Path, data: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        json.dumps(data, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


def build_payload(kind: str, files: list[Path], parser: Any) -> dict[str, Any]:
    parsed = [parser(path) for path in files]
    return {
        "format": kind,
        "file_count": len(parsed),
        "files": parsed,
    }


def main() -> None:
    project_root = Path(__file__).resolve().parent.parent
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--template-dir", type=Path, default=project_root / "template")
    parser.add_argument("--output-dir", type=Path, default=project_root / "output")
    args = parser.parse_args()

    template_dir = args.template_dir.resolve()
    output_dir = args.output_dir.resolve()
    docx_files = sorted(template_dir.glob("*.docx"))
    pptx_files = sorted(template_dir.glob("*.pptx"))
    if not docx_files:
        raise SystemExit(f"No DOCX files found in {template_dir}")
    if not pptx_files:
        raise SystemExit(f"No PPTX files found in {template_dir}")

    doc_payload = build_payload(
        "docx",
        docx_files,
        lambda path: parse_docx(path, output_dir),
    )
    ppt_payload = build_payload("pptx", pptx_files, parse_pptx)
    write_json(output_dir / "doc_structure.json", doc_payload)
    write_json(output_dir / "ppt_layout.json", ppt_payload)
    print(f"Wrote {output_dir / 'doc_structure.json'}")
    print(f"Wrote {output_dir / 'ppt_layout.json'}")


if __name__ == "__main__":
    main()
