#!/usr/bin/env python3
"""Validate rendered slide PNGs against editable PPT geometry."""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from typing import Any

from PIL import Image, ImageChops, ImageStat
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE


CHECKED_NAMES = (
    "title",
    "content",
    "chart_placeholder",
    "table_placeholder",
    "source",
    "page_number",
)


def normalize_png_names(render_dir: Path) -> dict[int, Path]:
    numbered: dict[int, Path] = {}
    for path in render_dir.glob("*.png"):
        match = re.search(r"(\d+)", path.stem)
        if not match:
            continue
        slide_number = int(match.group(1))
        target = render_dir / f"slide-{slide_number:02d}.png"
        if path != target:
            if target.exists():
                target.unlink()
            path.rename(target)
        numbered[slide_number] = target
    return numbered


def inside_slide(shape: Any, slide_width: int, slide_height: int) -> bool:
    return (
        shape.left >= 0
        and shape.top >= 0
        and shape.left + shape.width <= slide_width
        and shape.top + shape.height <= slide_height
    )


def overlap_ratio(a: Any, b: Any) -> float:
    width = max(
        0,
        min(a.left + a.width, b.left + b.width) - max(a.left, b.left),
    )
    height = max(
        0,
        min(a.top + a.height, b.top + b.height) - max(a.top, b.top),
    )
    intersection = width * height
    minimum_area = min(a.width * a.height, b.width * b.height)
    return intersection / minimum_area if minimum_area else 0.0


def pixel_box(shape: Any, slide_width: int, slide_height: int, image: Image.Image) -> tuple[int, int, int, int]:
    scale_x = image.width / slide_width
    scale_y = image.height / slide_height
    return (
        max(0, round(shape.left * scale_x)),
        max(0, round(shape.top * scale_y)),
        min(image.width, round((shape.left + shape.width) * scale_x)),
        min(image.height, round((shape.top + shape.height) * scale_y)),
    )


def background_color(image: Image.Image) -> tuple[int, int, int]:
    points = [
        image.getpixel((2, 2)),
        image.getpixel((image.width - 3, 2)),
        image.getpixel((2, image.height - 3)),
        image.getpixel((image.width - 3, image.height - 3)),
    ]
    return tuple(sorted(point[channel] for point in points)[len(points) // 2] for channel in range(3))


def edge_ink_ratio(
    image: Image.Image,
    box: tuple[int, int, int, int],
    background: tuple[int, int, int],
    edge: str,
    band: int = 3,
) -> float:
    left, top, right, bottom = box
    if right <= left or bottom <= top:
        return 1.0
    crop = image.crop(box).convert("RGB")
    bg = Image.new("RGB", crop.size, background)
    difference = ImageChops.difference(crop, bg).convert("L")
    mask = difference.point(lambda value: 255 if value > 28 else 0)
    if edge == "bottom":
        edge_crop = mask.crop((0, max(0, mask.height - band), mask.width, mask.height))
    elif edge == "right":
        edge_crop = mask.crop((max(0, mask.width - band), 0, mask.width, mask.height))
    else:
        raise ValueError(edge)
    stat = ImageStat.Stat(edge_crop)
    return round((stat.mean[0] / 255), 4)


def non_background_ratio(image: Image.Image, background: tuple[int, int, int]) -> float:
    bg = Image.new("RGB", image.size, background)
    difference = ImageChops.difference(image.convert("RGB"), bg).convert("L")
    mask = difference.point(lambda value: 255 if value > 20 else 0)
    return round(ImageStat.Stat(mask).mean[0] / 255, 4)


def check_text_shape(
    shape: Any,
    image: Image.Image,
    slide_width: int,
    slide_height: int,
    background: tuple[int, int, int],
) -> dict[str, Any]:
    box = pixel_box(shape, slide_width, slide_height, image)
    bottom_ratio = edge_ink_ratio(image, box, background, "bottom")
    right_ratio = edge_ink_ratio(image, box, background, "right")
    auto_fit = bool(
        getattr(shape, "has_text_frame", False)
        and shape.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    )
    overflow_failed = bottom_ratio > 0.20 or right_ratio > 0.20
    overflow_suspected = bottom_ratio > 0.08 or right_ratio > 0.08
    status = "fail" if overflow_failed else "warning" if overflow_suspected else "pass"
    return {
        "status": status,
        "inside_slide": inside_slide(shape, slide_width, slide_height),
        "powerpoint_auto_fit": auto_fit,
        "bottom_edge_ink_ratio": bottom_ratio,
        "right_edge_ink_ratio": right_ratio,
        "details": (
            "Rendered text strongly touches a box edge; clipping/hidden bullet is likely."
            if overflow_failed
            else "Rendered text touches a box edge; inspect for clipping."
            if overflow_suspected
            else "No rendered edge clipping detected."
        ),
    }


def build_report(
    pptx_path: Path,
    render_dir: Path,
    plan_path: Path,
) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    plan = json.loads(plan_path.read_text(encoding="utf-8"))
    page_types = {
        slide["slide_number"]: slide["page_type"] for slide in plan["slides"]
    }
    pngs = normalize_png_names(render_dir)
    slide_reports: list[dict[str, Any]] = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        image_path = pngs.get(slide_number)
        if image_path is None:
            slide_reports.append(
                {
                    "slide_number": slide_number,
                    "page_type": page_types.get(slide_number, "unknown"),
                    "image": None,
                    "status": "fail",
                    "issues": ["Missing rendered PNG."],
                }
            )
            continue

        with Image.open(image_path) as opened:
            image = opened.convert("RGB")
        background = background_color(image)
        shapes = {shape.name: shape for shape in slide.shapes}
        checks: dict[str, Any] = {}
        issues: list[str] = []

        image_ratio = non_background_ratio(image, background)
        checks["image_render"] = {
            "status": "pass" if image.width > 0 and image.height > 0 and image_ratio > 0.002 else "fail",
            "width": image.width,
            "height": image.height,
            "non_background_ratio": image_ratio,
        }

        for shape_name, check_name in (
            ("title", "title_bounds"),
            ("content", "body_overflow"),
        ):
            shape = shapes.get(shape_name)
            if shape is None:
                checks[check_name] = {"status": "fail", "details": f"Missing {shape_name} shape."}
                issues.append(f"Missing {shape_name} shape.")
                continue
            checks[check_name] = check_text_shape(
                shape,
                image,
                prs.slide_width,
                prs.slide_height,
                background,
            )
            if checks[check_name]["status"] != "pass":
                issues.append(checks[check_name]["details"])

        table = shapes.get("table_placeholder")
        checks["table_region"] = {
            "status": (
                "pass"
                if table is None or inside_slide(table, prs.slide_width, prs.slide_height)
                else "fail"
            ),
            "present": table is not None,
            "inside_slide": (
                inside_slide(table, prs.slide_width, prs.slide_height)
                if table is not None
                else None
            ),
            "editable_native_table": bool(table is not None and table.has_table),
        }
        if checks["table_region"]["status"] == "fail":
            issues.append("Table exceeds slide bounds.")

        chart = shapes.get("chart_placeholder")
        chart_text = chart.text if chart is not None and chart.has_text_frame else ""
        checks["chart_placeholder"] = {
            "status": (
                "pass"
                if chart is None
                or (
                    inside_slide(chart, prs.slide_width, prs.slide_height)
                    and "图表区域" in chart_text
                )
                else "fail"
            ),
            "present": chart is not None,
            "inside_slide": (
                inside_slide(chart, prs.slide_width, prs.slide_height)
                if chart is not None
                else None
            ),
            "label_correct": "图表区域" in chart_text if chart is not None else None,
        }
        if checks["chart_placeholder"]["status"] == "fail":
            issues.append("Chart placeholder is outside bounds or incorrectly labeled.")

        overlap_pairs: list[dict[str, Any]] = []
        checked = [name for name in CHECKED_NAMES if name in shapes]
        for index, left_name in enumerate(checked):
            for right_name in checked[index + 1 :]:
                ratio = overlap_ratio(shapes[left_name], shapes[right_name])
                if ratio > 0.02:
                    overlap_pairs.append(
                        {
                            "shape_a": left_name,
                            "shape_b": right_name,
                            "intersection_ratio": round(ratio, 4),
                        }
                    )
        checks["overlap"] = {
            "status": "fail" if overlap_pairs else "pass",
            "pairs": overlap_pairs,
        }
        if overlap_pairs:
            issues.append(f"Detected {len(overlap_pairs)} unintended overlap(s).")

        statuses = [check["status"] for check in checks.values()]
        status = "fail" if "fail" in statuses else "warning" if "warning" in statuses else "pass"
        slide_reports.append(
            {
                "slide_number": slide_number,
                "page_type": page_types.get(slide_number, "unknown"),
                "image": image_path.name,
                "status": status,
                "checks": checks,
                "issues": issues,
            }
        )

    counts = {
        status: sum(slide["status"] == status for slide in slide_reports)
        for status in ("pass", "warning", "fail")
    }
    return {
        "source_pptx": pptx_path.name,
        "renderer": "Microsoft PowerPoint PNG export",
        "slide_count": len(prs.slides),
        "rendered_png_count": len(pngs),
        "overall_status": (
            "fail"
            if counts["fail"]
            else "warning"
            if counts["warning"]
            else "pass"
        ),
        "summary": counts,
        "checks_performed": [
            "title bounds and rendered edge clipping",
            "body bounds, auto-fit, and rendered edge clipping",
            "table bounds and native editability",
            "chart placeholder bounds and label",
            "pairwise element overlap",
            "PNG dimensions and non-blank content",
        ],
        "slides": slide_reports,
    }


def main() -> None:
    project_root = Path(__file__).resolve().parent.parent
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--pptx",
        type=Path,
        default=project_root / "output" / "generated_test.pptx",
    )
    parser.add_argument(
        "--render-dir",
        type=Path,
        default=project_root / "output" / "render_check",
    )
    parser.add_argument(
        "--plan",
        type=Path,
        default=project_root / "output" / "presentation_plan.json",
    )
    parser.add_argument(
        "--report",
        type=Path,
        default=project_root / "output" / "render_check" / "check_report.json",
    )
    args = parser.parse_args()

    args.render_dir.mkdir(parents=True, exist_ok=True)
    report = build_report(args.pptx, args.render_dir, args.plan)
    args.report.write_text(
        json.dumps(report, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(
        f"Wrote {args.report}: {report['summary']} "
        f"(overall={report['overall_status']})"
    )


if __name__ == "__main__":
    main()
