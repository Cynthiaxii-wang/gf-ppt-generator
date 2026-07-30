#!/usr/bin/env python3
"""Extract reusable visual style tokens from a PowerPoint template."""

from __future__ import annotations

import argparse
import json
import statistics
import zipfile
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any, Iterable
from xml.etree import ElementTree as ET

import yaml
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE


EMU_PER_INCH = 914400
NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
ROLE_KEYS = (
    "cover_title",
    "section_title",
    "body_title",
    "body",
    "table",
    "chart",
    "source",
    "footer",
)


def inches(value: int | float) -> float:
    return round(float(value) / EMU_PER_INCH, 3)


def mode_or_unknown(values: Iterable[Any]) -> Any:
    cleaned = [value for value in values if value not in (None, "", "unknown")]
    return Counter(cleaned).most_common(1)[0][0] if cleaned else "unknown"


def median_or_unknown(values: Iterable[int | float]) -> Any:
    cleaned = [float(value) for value in values if value is not None]
    return round(statistics.median(cleaned), 3) if cleaned else "unknown"


def load_mapping(path: Path) -> dict[int, dict[str, Any]]:
    payload = yaml.safe_load(path.read_text(encoding="utf-8"))
    return {item["slide_number"]: item for item in payload["slides"]}


def iter_shapes(shapes: Iterable[Any], prefix: str = "") -> Iterable[tuple[str, Any]]:
    for shape in shapes:
        path = f"{prefix}/{shape.name}" if prefix else shape.name
        yield path, shape
        if hasattr(shape, "shapes"):
            yield from iter_shapes(shape.shapes, path)


def iter_layout_shapes(shapes: list[dict[str, Any]], prefix: str = "") -> Iterable[tuple[str, dict[str, Any]]]:
    for shape in shapes:
        path = f"{prefix}/{shape['name']}" if prefix else shape["name"]
        yield path, shape
        yield from iter_layout_shapes(shape.get("children", []), path)


def find_shape(slide: Any, configured_name: Any) -> Any | None:
    if not configured_name:
        return None
    name = configured_name[0] if isinstance(configured_name, list) else configured_name
    if "/" in name:
        for path, shape in iter_shapes(slide.shapes):
            if path == name:
                return shape
        return None
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def color_value(color: Any) -> str | None:
    try:
        if color.type == MSO_COLOR_TYPE.RGB and color.rgb is not None:
            return f"#{color.rgb}"
        if color.type == MSO_COLOR_TYPE.SCHEME and color.theme_color is not None:
            return f"theme:{color.theme_color}"
    except (AttributeError, ValueError):
        return None
    return None


def xml_text_properties(shape: Any) -> dict[str, list[Any]]:
    values: dict[str, list[Any]] = {"fonts": [], "sizes": [], "colors": []}
    if not hasattr(shape, "_element"):
        return values
    for prop in shape._element.xpath(".//a:rPr | .//a:defRPr | .//a:endParaRPr"):
        size = prop.get("sz")
        if size and size.isdigit():
            values["sizes"].append(int(size) / 100)
        for font_tag in ("a:latin", "a:ea"):
            font_node = prop.find(font_tag, prop.nsmap)
            if font_node is not None and font_node.get("typeface"):
                typeface = font_node.get("typeface")
                if not typeface.startswith("+"):
                    values["fonts"].append(typeface)
        solid = prop.find("a:solidFill", prop.nsmap)
        if solid is not None:
            rgb = solid.find("a:srgbClr", prop.nsmap)
            scheme = solid.find("a:schemeClr", prop.nsmap)
            if rgb is not None and rgb.get("val"):
                values["colors"].append(f"#{rgb.get('val')}")
            elif scheme is not None and scheme.get("val"):
                values["colors"].append(f"theme:{scheme.get('val')}")
    return values


def collect_text_style(shape: Any) -> dict[str, list[Any]]:
    values: dict[str, list[Any]] = {
        "fonts": [],
        "sizes": [],
        "colors": [],
        "bold": [],
        "line_spacing": [],
        "space_before": [],
        "space_after": [],
    }
    if not getattr(shape, "has_text_frame", False):
        return values
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.line_spacing is not None:
            if hasattr(paragraph.line_spacing, "pt"):
                values["line_spacing"].append(paragraph.line_spacing.pt)
            else:
                values["line_spacing"].append(float(paragraph.line_spacing))
        if paragraph.space_before is not None:
            values["space_before"].append(paragraph.space_before.pt)
        if paragraph.space_after is not None:
            values["space_after"].append(paragraph.space_after.pt)
        for run in paragraph.runs:
            if not run.text.strip():
                continue
            if run.font.name:
                values["fonts"].append(run.font.name)
            if run.font.size:
                values["sizes"].append(run.font.size.pt)
            run_color = color_value(run.font.color)
            if run_color:
                values["colors"].append(run_color)
            if run.font.bold is not None:
                values["bold"].append(bool(run.font.bold))
    xml_values = xml_text_properties(shape)
    for key in ("fonts", "sizes", "colors"):
        values[key].extend(xml_values[key])
    return values


def collect_inherited_text_style(slide: Any, shape: Any) -> dict[str, list[Any]]:
    values: dict[str, list[Any]] = defaultdict(list)
    if not getattr(shape, "is_placeholder", False):
        return values
    placeholder_index = shape.placeholder_format.idx
    for placeholder in slide.slide_layout.placeholders:
        if placeholder.placeholder_format.idx == placeholder_index:
            collected = collect_text_style(placeholder)
            for key, entries in collected.items():
                values[key].extend(entries)
            break
    for placeholder in slide.slide_layout.slide_master.placeholders:
        if placeholder.placeholder_format.idx == placeholder_index:
            collected = collect_text_style(placeholder)
            for key, entries in collected.items():
                values[key].extend(entries)
            break
    return values


def collect_table_style(shape: Any) -> tuple[dict[str, list[Any]], list[str]]:
    values: dict[str, list[Any]] = defaultdict(list)
    fills: list[str] = []
    if not getattr(shape, "has_table", False):
        return values, fills
    for row in shape.table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue
                    if run.font.name:
                        values["fonts"].append(run.font.name)
                    if run.font.size:
                        values["sizes"].append(run.font.size.pt)
                    run_color = color_value(run.font.color)
                    if run_color:
                        values["colors"].append(run_color)
            try:
                fill_color = color_value(cell.fill.fore_color)
            except TypeError:
                fill_color = None
            if fill_color:
                fills.append(fill_color)
    return values, fills


def chart_xml_styles(pptx_path: Path) -> dict[str, list[Any]]:
    values: dict[str, list[Any]] = {"fonts": [], "sizes": [], "colors": []}
    with zipfile.ZipFile(pptx_path) as archive:
        for name in archive.namelist():
            if not name.startswith("ppt/charts/chart") or not name.endswith(".xml"):
                continue
            root = ET.fromstring(archive.read(name))
            for node in root.findall(".//a:latin", NS) + root.findall(".//a:ea", NS):
                typeface = node.get("typeface")
                if typeface and not typeface.startswith("+"):
                    values["fonts"].append(typeface)
            for node in root.findall(".//*[@sz]"):
                size = node.get("sz")
                if size and size.isdigit():
                    values["sizes"].append(int(size) / 100)
            for node in root.findall(".//a:srgbClr", NS):
                if node.get("val"):
                    values["colors"].append(f"#{node.get('val')}")
    return values


def theme_tokens(pptx_path: Path) -> dict[str, Any]:
    with zipfile.ZipFile(pptx_path) as archive:
        theme_name = next(
            name for name in archive.namelist() if name.startswith("ppt/theme/theme")
        )
        root = ET.fromstring(archive.read(theme_name))
    major = root.find(".//a:majorFont", NS)
    minor = root.find(".//a:minorFont", NS)

    def theme_font(node: ET.Element | None, kind: str) -> str:
        if node is None:
            return "unknown"
        hans = node.find("a:font[@script='Hans']", NS)
        if hans is not None and hans.get("typeface"):
            return hans.get("typeface", "unknown")
        latin = node.find("a:latin", NS)
        return latin.get("typeface", "unknown") if latin is not None else "unknown"

    colors: dict[str, str] = {}
    scheme = root.find(".//a:clrScheme", NS)
    if scheme is not None:
        for child in list(scheme):
            color_node = next(iter(child), None)
            if color_node is not None:
                value = color_node.get("lastClr") or color_node.get("val")
                if value:
                    colors[child.tag.rsplit("}", 1)[-1]] = f"#{value}"
    return {
        "major_hans": theme_font(major, "major"),
        "minor_hans": theme_font(minor, "minor"),
        "colors": colors,
    }


def region(boxes: list[dict[str, int]]) -> dict[str, Any]:
    if not boxes:
        return {
            "left_emu": "unknown",
            "top_emu": "unknown",
            "width_emu": "unknown",
            "height_emu": "unknown",
            "left_in": "unknown",
            "top_in": "unknown",
            "width_in": "unknown",
            "height_in": "unknown",
        }
    result: dict[str, Any] = {}
    for key in ("left", "top", "width", "height"):
        value = int(statistics.median(box[key] for box in boxes))
        result[f"{key}_emu"] = value
        result[f"{key}_in"] = inches(value)
    return result


def style_entry(
    values: dict[str, list[Any]],
    theme_fallback: str,
    prefer_cjk: bool = True,
) -> dict[str, Any]:
    font_counts = Counter(
        value for value in values.get("fonts", []) if value not in (None, "", "unknown")
    )
    candidates = [name for name, _ in font_counts.most_common(5)]
    cjk_markers = ("Han Sans", "思源", "等线", "黑体", "宋体", "ゴシック")
    font = next(
        (name for name in candidates if any(marker in name for marker in cjk_markers)),
        candidates[0] if candidates else "unknown",
    ) if prefer_cjk else (candidates[0] if candidates else "unknown")
    size = median_or_unknown(values.get("sizes", []))
    color = mode_or_unknown(values.get("colors", []))
    return {
        "font_family": font if font != "unknown" else theme_fallback,
        "font_candidates": candidates or [theme_fallback],
        "font_size_pt": size,
        "font_color": color,
        "bold": mode_or_unknown(values.get("bold", [])),
        "status": (
            "observed"
            if font != "unknown" and size != "unknown"
            else "partial"
            if font != "unknown" or size != "unknown" or theme_fallback != "unknown"
            else "unknown"
        ),
    }


def build_config(
    pptx_path: Path,
    layout_path: Path,
    mapping_path: Path,
) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    mapping = load_mapping(mapping_path)
    layout_payload = json.loads(layout_path.read_text(encoding="utf-8"))["files"][0]
    theme = theme_tokens(pptx_path)
    role_values: dict[str, dict[str, list[Any]]] = {
        role: defaultdict(list) for role in ROLE_KEYS
    }
    table_fills: list[str] = []
    emphasis_colors: list[str] = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        mapped = mapping.get(slide_number, {})
        title_shape = find_shape(slide, mapped.get("title_shape"))
        content_shape = find_shape(slide, mapped.get("content_shape"))
        page_type = mapped.get("page_type")
        title_role = (
            "cover_title"
            if page_type == "cover"
            else "section_title"
            if page_type == "section"
            else "body_title"
        )
        if title_shape is not None:
            collected = collect_text_style(title_shape)
            inherited = collect_inherited_text_style(slide, title_shape)
            for key, entries in inherited.items():
                collected[key].extend(entries)
            for key, entries in collected.items():
                role_values[title_role][key].extend(entries)
        if content_shape is not None:
            collected = collect_text_style(content_shape)
            inherited = collect_inherited_text_style(slide, content_shape)
            for key, entries in inherited.items():
                collected[key].extend(entries)
            for key, entries in collected.items():
                role_values["body"][key].extend(entries)
            if getattr(content_shape, "has_text_frame", False):
                for paragraph in content_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run_color = color_value(run.font.color)
                        font_name = run.font.name or ""
                        if (
                            run.text.strip()
                            and run_color
                            and (run.font.bold or "Bold" in font_name)
                        ):
                            emphasis_colors.append(run_color)

        for _, shape in iter_shapes(slide.shapes):
            if getattr(shape, "has_table", False):
                collected, fills = collect_table_style(shape)
                for key, entries in collected.items():
                    role_values["table"][key].extend(entries)
                table_fills.extend(fills)
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                role = None
                if text.startswith("数据来源"):
                    role = "source"
                elif "编号" in shape.name or (text.isdigit() and shape.top > prs.slide_height * 0.8):
                    role = "footer"
                if role:
                    collected = collect_text_style(shape)
                    for key, entries in collected.items():
                        role_values[role][key].extend(entries)

    chart_values = chart_xml_styles(pptx_path)
    for key, entries in chart_values.items():
        role_values["chart"][key].extend(entries)

    regions: dict[str, list[dict[str, int]]] = defaultdict(list)
    for slide in layout_payload["slides"]:
        mapped = mapping.get(slide["slide_number"], {})
        for path, shape in iter_layout_shapes(slide["shapes"]):
            name = shape["name"]
            configured_title = mapped.get("title_shape")
            configured_content = mapped.get("content_shape")
            configured_charts = mapped.get("chart_shape") or []
            configured_tables = mapped.get("table_shape") or []
            if isinstance(configured_title, str) and name == configured_title:
                regions["title"].append(shape)
            if isinstance(configured_content, str) and name == configured_content:
                regions["content"].append(shape)
            if path in configured_charts or name in configured_charts:
                regions["chart"].append(shape)
            if path in configured_tables or name in configured_tables:
                regions["table"].append(shape)
            if (shape.get("text") or "").strip().startswith("数据来源"):
                regions["source"].append(shape)
            if "编号" in name:
                regions["footer"].append(shape)

    title_region = region(regions["title"])
    content_region = region(regions["content"])
    left_candidates = [
        value
        for value in (
            title_region["left_emu"],
            content_region["left_emu"],
        )
        if isinstance(value, int)
    ]
    right_candidates = [
        prs.slide_width - (box["left"] + box["width"])
        for box in regions["title"] + regions["content"]
    ]
    title_gap_values = []
    for slide in layout_payload["slides"]:
        mapped = mapping.get(slide["slide_number"], {})
        title = next(
            (
                shape
                for _, shape in iter_layout_shapes(slide["shapes"])
                if shape["name"] == mapped.get("title_shape")
            ),
            None,
        )
        content = next(
            (
                shape
                for _, shape in iter_layout_shapes(slide["shapes"])
                if shape["name"] == mapped.get("content_shape")
            ),
            None,
        )
        if title and content and content["top"] >= title["top"] + title["height"]:
            title_gap_values.append(content["top"] - title["top"] - title["height"])

    entries = {
        role: style_entry(
            role_values[role],
            theme["major_hans"] if "title" in role else theme["minor_hans"],
            prefer_cjk=role not in {"chart", "source"},
        )
        for role in ROLE_KEYS
    }
    background_colors = Counter(
        color
        for color in (
            theme["colors"].get("lt1"),
            "#353160",
        )
        if color
    )

    return {
        "source_template": pptx_path.name,
        "style": {
            "font": {
                "title_style": entries["body_title"]["font_family"],
                "cover_title_style": entries["cover_title"]["font_family"],
                "section_title_style": entries["section_title"]["font_family"],
                "body_style": entries["body"]["font_family"],
                "table_style": entries["table"]["font_family"],
                "chart_style": entries["chart"]["font_family"],
                "source_style": entries["source"]["font_family"],
                "footer_style": entries["footer"]["font_family"],
                "theme_major_hans": theme["major_hans"],
                "theme_minor_hans": theme["minor_hans"],
            },
            "color": {
                "title_color": entries["body_title"]["font_color"],
                "body_color": entries["body"]["font_color"],
                "emphasis_color": mode_or_unknown(emphasis_colors),
                "background_color": mode_or_unknown(background_colors.elements()),
                "section_background_color": mode_or_unknown(background_colors.elements()),
                "section_accent_color": "#353160",
                "table_text_color": entries["table"]["font_color"],
                "table_fill_color": mode_or_unknown(table_fills),
                "theme_palette": theme["colors"],
            },
            "size": {
                "cover_title_pt": entries["cover_title"]["font_size_pt"],
                "section_title_pt": entries["section_title"]["font_size_pt"],
                "body_title_pt": entries["body_title"]["font_size_pt"],
                "body_text_pt": entries["body"]["font_size_pt"],
                "table_text_pt": entries["table"]["font_size_pt"],
                "chart_text_pt": entries["chart"]["font_size_pt"],
                "annotation_pt": entries["source"]["font_size_pt"],
                "page_number_pt": entries["footer"]["font_size_pt"],
            },
            "spacing": {
                "title_to_content_gap_emu": (
                    int(statistics.median(title_gap_values))
                    if title_gap_values
                    else "unknown"
                ),
                "title_to_content_gap_in": (
                    inches(statistics.median(title_gap_values))
                    if title_gap_values
                    else "unknown"
                ),
                "body_line_spacing": median_or_unknown(role_values["body"]["line_spacing"]),
                "body_space_before_pt": median_or_unknown(role_values["body"]["space_before"]),
                "body_space_after_pt": median_or_unknown(role_values["body"]["space_after"]),
            },
            "layout": {
                "slide_width_emu": int(prs.slide_width),
                "slide_height_emu": int(prs.slide_height),
                "slide_width_in": inches(prs.slide_width),
                "slide_height_in": inches(prs.slide_height),
                "page_margin_left_emu": int(min(left_candidates)) if left_candidates else "unknown",
                "page_margin_right_emu": int(min(right_candidates)) if right_candidates else "unknown",
                "page_margin_left_in": inches(min(left_candidates)) if left_candidates else "unknown",
                "page_margin_right_in": inches(min(right_candidates)) if right_candidates else "unknown",
                "title_region": title_region,
                "content_region": content_region,
                "chart_region": region(regions["chart"]),
                "table_region": region(regions["table"]),
                "source_region": region(regions["source"]),
                "footer_region": region(regions["footer"]),
            },
            "element_status": entries,
            "unknown_policy": "无法从模板直接样式、主题或可观测实例可靠解析的字段标记为unknown。",
        },
    }


def main() -> None:
    project_root = Path(__file__).resolve().parent.parent
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--pptx",
        type=Path,
        default=next((project_root / "template").glob("*.pptx")),
    )
    parser.add_argument(
        "--layout",
        type=Path,
        default=project_root / "output" / "ppt_layout.json",
    )
    parser.add_argument(
        "--mapping",
        type=Path,
        default=project_root / "output" / "slide_mapping.yaml",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=project_root / "output" / "style_config.yaml",
    )
    args = parser.parse_args()

    config = build_config(args.pptx, args.layout, args.mapping)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    args.output.write_text(
        yaml.safe_dump(
            config,
            allow_unicode=True,
            sort_keys=False,
            default_flow_style=False,
            width=110,
        ),
        encoding="utf-8",
    )
    print(f"Wrote {args.output}")


if __name__ == "__main__":
    main()
