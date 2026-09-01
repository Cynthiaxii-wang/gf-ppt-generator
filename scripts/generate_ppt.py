#!/usr/bin/env python3
"""Generate an editable structural PowerPoint from planning and style configs."""

from __future__ import annotations

import argparse
import json
import math
import re
import secrets
import threading
import uuid
from contextlib import contextmanager
from copy import deepcopy
from collections import defaultdict
from pathlib import Path
from typing import Any

import yaml
from PIL import Image
from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.opc.package import Part
from pptx.util import Emu, Pt


DEFAULT_COLOR = "#353160"
EMU_PER_INCH = 914400
EMU_PER_POINT = 12700
MIN_VISUAL_HEIGHT = 320000
REGION_GAP = 120000
MAX_EMBEDDED_IMAGE_BYTES = 100 * 1024 * 1024
MAX_EMBEDDED_IMAGE_PIXELS = 500_000_000
_PIL_IMAGE_LIMIT_LOCK = threading.RLock()

FIXED_DISCLAIMER_PARAGRAPHS = (
    "广发证券股份有限公司（以下简称“广发证券”）具备证券投资咨询业务资格。本报告只发送给广发证券重点客户，不对外公开发布，只有接收客户才可以使用，且对于接收客户而言具有相关保密义务。广发证券并不因相关人员通过其他途径收到或阅读本报告而视其为广发证券的客户。本报告的内容、观点或建议并未考虑个别客户的特定状况，不应被视为对特定客户关于特定证券或金融工具的投资建议。本报告发送给某客户是基于该客户被认为有能力独立评估投资风险、独立行使投资决策并独立承担相应风险。",
    "本报告所载资料的来源及观点的出处皆被广发证券股份有限公司认为可靠，但广发证券不对其准确性或完整性做出任何保证。报告内容仅供参考，报告中的信息或所表达观点不构成所涉证券买卖的出价或询价。广发证券不对因使用本报告的内容而引致的损失承担任何责任，除非法律法规有明确规定。客户不应以本报告取代其独立判断或仅根据本报告做出决策。",
    "广发证券可发出其它与本报告所载信息不一致及有不同结论的报告。本报告反映研究人员的不同观点、见解及分析方法，并不代表广发证券或其附属机构的立场。报告所载资料、意见及推测仅反映研究人员于发出本报告当日的判断，可随时更改且不予通告。",
    "本报告旨在发送给广发证券的特定客户及其它专业人士。未经广发证券事先书面许可，任何机构或个人不得以任何形式翻版、复制、刊登、转载和引用，否则由此造成的一切不良后果及法律责任由私自翻版、复制、刊登、转载和引用者承担。",
)


class EmbeddedImageError(ValueError):
    """An extracted Word image cannot be embedded safely."""


@contextmanager
def allow_large_image_metadata():
    """Temporarily bypass Pillow's pixel warning while reading image headers.

    The lock keeps Pillow's process-global setting isolated across concurrent
    Streamlit sessions.  Callers must still enforce the explicit hard limits
    above and must not decode the full raster while this context is active.
    """

    with _PIL_IMAGE_LIMIT_LOCK:
        previous_limit = Image.MAX_IMAGE_PIXELS
        Image.MAX_IMAGE_PIXELS = None
        try:
            yield
        finally:
            Image.MAX_IMAGE_PIXELS = previous_limit


def rgb(value: Any, fallback: str = DEFAULT_COLOR) -> RGBColor:
    text = value if isinstance(value, str) and value.startswith("#") else fallback
    text = text.lstrip("#")
    return RGBColor.from_string(text[:6])


def numeric(value: Any, fallback: float) -> float:
    return float(value) if isinstance(value, (int, float)) else fallback


def set_run_fonts(
    run: Any,
    chinese_font: str,
    latin_font: str = "Arial",
) -> None:
    """Assign the correct primary font to a single-script PowerPoint run."""

    drawing_namespace = (
        "http://schemas.openxmlformats.org/drawingml/2006/main"
    )
    chinese_font = chinese_font.replace(
        "Source Han Sans CN",
        "思源黑体 CN",
    )
    is_chinese = bool(
        re.search(r"[\u2e80-\u9fff\uf900-\ufaff\uff00-\uffef]", run.text)
    )
    primary_font = chinese_font if is_chinese else latin_font
    run.font.name = primary_font
    properties = run._r.get_or_add_rPr()
    latin = properties.find(f"{{{drawing_namespace}}}latin")
    if latin is None:
        latin = etree.SubElement(
            properties,
            f"{{{drawing_namespace}}}latin",
        )
    latin.set("typeface", primary_font)
    east_asian = properties.find(f"{{{drawing_namespace}}}ea")
    if east_asian is None:
        east_asian = etree.Element(f"{{{drawing_namespace}}}ea")
        latin.addnext(east_asian)
    east_asian.set("typeface", chinese_font)
    properties.set("lang", "zh-CN" if is_chinese else "en-US")


def script_segments(text: str) -> list[str]:
    """Split mixed Chinese/Latin text so PowerPoint can show each font."""

    return re.findall(
        r"[\u2e80-\u9fff\uf900-\ufaff\uff00-\uffef]+"
        r"|[^\u2e80-\u9fff\uf900-\ufaff\uff00-\uffef]+",
        text,
    ) or [""]


def load_inputs(output_dir: Path) -> dict[str, Any]:
    files = {
        "plan": output_dir / "presentation_plan.json",
        "content": output_dir / "presentation_content.json",
        "mapping": output_dir / "slide_mapping.yaml",
        "layout": output_dir / "ppt_layout.json",
        "style": output_dir / "style_config.yaml",
        "doc": output_dir / "doc_structure.json",
    }
    missing = [str(path) for path in files.values() if not path.exists()]
    if missing:
        raise FileNotFoundError(f"Missing required inputs: {missing}")
    style = yaml.safe_load(files["style"].read_text(encoding="utf-8"))
    source_template = style.get("source_template")
    return {
        "plan": json.loads(files["plan"].read_text(encoding="utf-8")),
        "content": json.loads(files["content"].read_text(encoding="utf-8")),
        "mapping": yaml.safe_load(files["mapping"].read_text(encoding="utf-8")),
        "layout": json.loads(files["layout"].read_text(encoding="utf-8")),
        "style": style,
        "doc": json.loads(files["doc"].read_text(encoding="utf-8")),
        "template_path": (
            output_dir.parent / "template" / source_template
            if source_template
            else None
        ),
        "asset_root": output_dir.parent,
    }


def validate_inputs(data: dict[str, Any]) -> None:
    planned = data["plan"]["slides"]
    content = data["content"]["slides"]
    if data["plan"]["planned_slide_count"] != len(planned):
        raise ValueError("presentation_plan.json slide count is inconsistent")
    if len(planned) != len(content):
        raise ValueError(
            f"Plan/content mismatch: {len(planned)} planned slides vs {len(content)} content slides"
        )
    if [item["slide_number"] for item in planned] != [
        item["slide_number"] for item in content
    ]:
        raise ValueError("Plan/content slide numbers are not aligned")


def iter_layout_shapes(shapes: list[dict[str, Any]], prefix: str = ""):
    for shape in shapes:
        path = f"{prefix}/{shape['name']}" if prefix else shape["name"]
        yield path, shape
        yield from iter_layout_shapes(shape.get("children", []), path)


def shape_by_name(slide_layout: dict[str, Any], configured: Any) -> dict[str, Any] | None:
    if not configured:
        return None
    names = configured if isinstance(configured, list) else [configured]
    for wanted in names:
        for path, shape in iter_layout_shapes(slide_layout["shapes"]):
            if path == wanted or shape["name"] == wanted:
                if path == wanted and "/" in wanted:
                    parent_name = wanted.split("/", 1)[0]
                    parent = next(
                        (
                            top_level
                            for top_level in slide_layout["shapes"]
                            if top_level["name"] == parent_name
                        ),
                        None,
                    )
                    if parent is not None:
                        return parent
                return shape
    return None


def shape_by_text_prefix(slide_layout: dict[str, Any], prefix: str) -> dict[str, Any] | None:
    return next(
        (
            shape
            for _, shape in iter_layout_shapes(slide_layout["shapes"])
            if (shape.get("text") or "").strip().startswith(prefix)
        ),
        None,
    )


def shape_by_name_fragment(slide_layout: dict[str, Any], fragment: str) -> dict[str, Any] | None:
    return next(
        (
            shape
            for _, shape in iter_layout_shapes(slide_layout["shapes"])
            if fragment in shape["name"]
        ),
        None,
    )


def region_from_style(style: dict[str, Any], key: str) -> dict[str, int]:
    region = style["style"]["layout"][key]
    fallback = {
        "left_emu": 685800,
        "top_emu": 1363926,
        "width_emu": 10831195,
        "height_emu": 957700,
    }
    return {
        "left": int(region.get("left_emu", fallback["left_emu"]))
        if isinstance(region.get("left_emu"), (int, float))
        else fallback["left_emu"],
        "top": int(region.get("top_emu", fallback["top_emu"]))
        if isinstance(region.get("top_emu"), (int, float))
        else fallback["top_emu"],
        "width": int(region.get("width_emu", fallback["width_emu"]))
        if isinstance(region.get("width_emu"), (int, float))
        else fallback["width_emu"],
        "height": int(region.get("height_emu", fallback["height_emu"]))
        if isinstance(region.get("height_emu"), (int, float))
        else fallback["height_emu"],
    }


def geometry(shape: dict[str, Any] | None, fallback: dict[str, int]) -> dict[str, int]:
    if shape is None:
        return fallback
    return {
        "left": int(shape["left"]),
        "top": int(shape["top"]),
        "width": int(shape["width"]),
        "height": int(shape["height"]),
    }


def fit_vertical(
    box: dict[str, int],
    minimum_top: int,
    maximum_bottom: int,
    minimum_height: int = 300000,
) -> dict[str, int]:
    fitted = dict(box)
    fitted["top"] = max(fitted["top"], minimum_top)
    fitted["height"] = min(fitted["height"], maximum_bottom - fitted["top"])
    if fitted["height"] < minimum_height:
        fitted["top"] = minimum_top
        fitted["height"] = max(minimum_height, maximum_bottom - minimum_top)
    return fitted


def box_bottom(box: dict[str, int]) -> int:
    return box["top"] + box["height"]


def box_right(box: dict[str, int]) -> int:
    return box["left"] + box["width"]


def boxes_overlap(first: dict[str, int], second: dict[str, int]) -> bool:
    return not (
        box_right(first) <= second["left"]
        or box_right(second) <= first["left"]
        or box_bottom(first) <= second["top"]
        or box_bottom(second) <= first["top"]
    )


def box_debug(box: dict[str, int]) -> dict[str, Any]:
    return {
        **box,
        "left_in": round(box["left"] / EMU_PER_INCH, 3),
        "top_in": round(box["top"] / EMU_PER_INCH, 3),
        "width_in": round(box["width"] / EMU_PER_INCH, 3),
        "height_in": round(box["height"] / EMU_PER_INCH, 3),
    }


def estimate_text_height(
    text: str,
    width: int,
    font_size_pt: float,
    *,
    bullet: bool = False,
    margin: int = 45720,
    paragraph_space_pt: float = 5,
) -> tuple[int, int]:
    """Estimate wrapped text height without relying on a PowerPoint renderer."""
    usable_width = max(width - 2 * margin, 1)
    average_character_width = max(font_size_pt * EMU_PER_POINT * 0.92, 1)
    characters_per_line = max(1, int(usable_width / average_character_width))
    line_count = 0
    paragraphs = text.splitlines() or [""]
    for paragraph in paragraphs:
        character_count = len(paragraph) + (2 if bullet else 0)
        line_count += max(1, math.ceil(character_count / characters_per_line))
    line_height = font_size_pt * EMU_PER_POINT * 1.28
    paragraph_spacing = max(0, len(paragraphs) - 1) * paragraph_space_pt * EMU_PER_POINT
    height = int(line_count * line_height + paragraph_spacing + 2 * margin)
    return height, line_count


def fit_title_font_size(
    text: str,
    box: dict[str, int],
    preferred_size_pt: float,
    minimum_size_pt: float = 8,
) -> tuple[float, int]:
    """Shrink a normal slide title until it fits its configured title box."""
    size = preferred_size_pt
    while size > minimum_size_pt:
        height, lines = estimate_text_height(
            text,
            box["width"],
            size,
            margin=0,
            paragraph_space_pt=0,
        )
        if height <= box["height"]:
            return size, lines
        size = max(minimum_size_pt, size - 0.5)
    height, lines = estimate_text_height(
        text,
        box["width"],
        minimum_size_pt,
        margin=0,
        paragraph_space_pt=0,
    )
    if height <= box["height"]:
        return minimum_size_pt, lines
    fitted_size = max(1.0, minimum_size_pt * box["height"] / height)
    _, fitted_lines = estimate_text_height(
        text,
        box["width"],
        fitted_size,
        margin=0,
        paragraph_space_pt=0,
    )
    return fitted_size, fitted_lines


def title_style_for_page(
    page_type: str,
    fonts: dict[str, Any],
    sizes: dict[str, Any],
) -> tuple[str, float]:
    if page_type == "cover":
        return fonts["cover_title_style"], numeric(sizes["cover_title_pt"], 38)
    if page_type == "section":
        return fonts["section_title_style"], numeric(sizes["section_title_pt"], 40)
    return fonts["title_style"], numeric(sizes["body_title_pt"], 32)


def dynamic_layout(
    *,
    page_type: str,
    title: str,
    body_text: str,
    title_seed: dict[str, int],
    content_seed: dict[str, int],
    source_box: dict[str, int],
    footer_box: dict[str, int],
    slide_width: int,
    title_size: float,
    body_size: float,
    has_chart: bool,
    has_table: bool,
    expand_title_height: bool,
) -> tuple[dict[str, dict[str, int]], dict[str, Any]]:
    """Allocate non-overlapping regions, giving text priority over visuals."""
    title_height, title_lines = estimate_text_height(
        title,
        title_seed["width"],
        title_size,
        margin=0,
        paragraph_space_pt=0,
    )
    title_box = dict(title_seed)
    if expand_title_height:
        title_box["height"] = max(title_seed["height"], title_height)

    content_top = box_bottom(title_box) + REGION_GAP
    bottom_limit = min(source_box["top"], footer_box["top"]) - REGION_GAP
    available_height = max(0, bottom_limit - content_top)
    body_height, body_lines = estimate_text_height(
        body_text,
        content_seed["width"],
        body_size,
        bullet=bool(body_text),
    )
    if not body_text:
        body_height = 0
        body_lines = 0

    has_visual = has_chart or has_table
    visual_gap = REGION_GAP if body_text and has_visual else 0
    visual_minimum = MIN_VISUAL_HEIGHT if has_visual else 0
    requested_body_height = min(body_height, available_height)
    if has_visual:
        body_budget = max(0, available_height - visual_gap - visual_minimum)
        actual_body_height = min(requested_body_height, body_budget)
    else:
        actual_body_height = requested_body_height

    content_box = {
        "left": content_seed["left"],
        "top": content_top,
        "width": content_seed["width"],
        "height": actual_body_height,
    }
    visual_top = content_top + actual_body_height + visual_gap
    visual_height = max(0, bottom_limit - visual_top)

    regions: dict[str, dict[str, int]] = {
        "title": title_box,
        "source": source_box,
        "footer": footer_box,
    }
    if body_text:
        regions["content"] = content_box

    page_left = min(title_seed["left"], content_seed["left"])
    page_right = max(box_right(title_seed), box_right(content_seed))
    visual_width = max(0, min(slide_width, page_right) - page_left)
    if has_chart and has_table:
        column_gap = REGION_GAP
        chart_width = max(0, (visual_width - column_gap) // 2)
        regions["chart"] = {
            "left": page_left,
            "top": visual_top,
            "width": chart_width,
            "height": visual_height,
        }
        regions["table"] = {
            "left": page_left + chart_width + column_gap,
            "top": visual_top,
            "width": max(0, visual_width - chart_width - column_gap),
            "height": visual_height,
        }
    elif has_chart:
        regions["chart"] = {
            "left": page_left,
            "top": visual_top,
            "width": visual_width,
            "height": visual_height,
        }
    elif has_table:
        regions["table"] = {
            "left": page_left,
            "top": visual_top,
            "width": visual_width,
            "height": visual_height,
        }

    visible_regions = {
        key: value
        for key, value in regions.items()
        if value["width"] > 0 and value["height"] > 0
    }
    overlaps = []
    region_items = list(visible_regions.items())
    for index, (first_name, first_box) in enumerate(region_items):
        for second_name, second_box in region_items[index + 1 :]:
            if boxes_overlap(first_box, second_box):
                overlaps.append([first_name, second_name])

    debug = {
        "page_type": page_type,
        "title_estimated_lines": title_lines,
        "body_estimated_lines": body_lines,
        "body_required_height_emu": body_height,
        "body_allocated_height_emu": actual_body_height,
        "body_auto_fit_required": body_height > actual_body_height,
        "available_height_emu": available_height,
        "visual_scaled": has_visual,
        "regions": {name: box_debug(box) for name, box in regions.items()},
        "overlaps": overlaps,
    }
    return regions, debug


def add_text_box(
    slide: Any,
    box: dict[str, int],
    text: str,
    name: str,
    font_name: str,
    font_size: float,
    color: RGBColor,
    bold: bool = False,
    bullet: bool = False,
    margin: int = 45720,
) -> Any:
    shape = slide.shapes.add_textbox(
        Emu(box["left"]),
        Emu(box["top"]),
        Emu(box["width"]),
        Emu(box["height"]),
    )
    shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Emu(margin)
    frame.margin_right = Emu(margin)
    frame.margin_top = Emu(margin)
    frame.margin_bottom = Emu(margin)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    lines = text.splitlines() or [""]
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.alignment = PP_ALIGN.LEFT
        if bullet:
            paragraph.text = f"• {line}"
        paragraph.space_after = Pt(5)
        paragraph.line_spacing = 1.08
        for run in paragraph.runs:
            set_run_fonts(run, font_name)
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold
    return shape


def add_chart_placeholder(
    slide: Any,
    box: dict[str, int],
    font_name: str,
    font_size: float,
    title_color: RGBColor,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Emu(box["left"]),
        Emu(box["top"]),
        Emu(box["width"]),
        Emu(box["height"]),
    )
    shape.name = "chart_placeholder"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 245, 250)
    shape.line.color.rgb = RGBColor(185, 190, 229)
    shape.line.width = Pt(1)
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = "图表区域（待生成）"
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.runs[0]
    set_run_fonts(run, font_name)
    run.font.size = Pt(font_size)
    run.font.color.rgb = title_color


def table_cell_style(
    style_values: dict[str, Any],
    region: str,
    *,
    fallback_font: str,
    fallback_size: float,
    fallback_text_color: RGBColor,
    fallback_fill_color: RGBColor | None,
    fallback_bold: bool,
    fallback_alignment: PP_ALIGN,
) -> dict[str, Any]:
    table_config = style_values.get("table", style_values.get("table_style", {}))
    configured = table_config.get(region, {}) if isinstance(table_config, dict) else {}
    font_name = configured.get(
        "font_name",
        configured.get("font_family", configured.get("font", fallback_font)),
    )
    size = numeric(
        configured.get("font_size_pt", configured.get("size_pt")),
        fallback_size,
    )
    color_value = configured.get(
        "text_color",
        configured.get("font_color"),
    )
    fill_value = configured.get(
        "fill_color",
        configured.get("background_color"),
    )
    alignment_name = str(configured.get("alignment", "")).lower()
    alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(alignment_name, fallback_alignment)
    return {
        "font_name": font_name,
        "font_size": size,
        "text_color": rgb(color_value, str(fallback_text_color))
        if color_value
        else fallback_text_color,
        "fill_color": rgb(fill_value, "#FFFFFF")
        if fill_value
        else fallback_fill_color,
        "bold": bool(configured.get("bold", fallback_bold)),
        "alignment": alignment,
    }


def add_table_placeholder(
    slide: Any,
    box: dict[str, int],
    styles: dict[str, dict[str, Any]],
) -> None:
    table_shape = slide.shapes.add_table(
        4,
        3,
        Emu(box["left"]),
        Emu(box["top"]),
        Emu(box["width"]),
        Emu(box["height"]),
    )
    table_shape.name = "table_placeholder"
    table = table_shape.table
    labels = (
        ("字段", "数值", "说明"),
        ("待填充", "—", "根据原文表格生成"),
        ("待填充", "—", "保持关键数据"),
        ("待填充", "—", "保持比较关系"),
    )
    for row_index, row in enumerate(table.rows):
        for column_index, cell in enumerate(row.cells):
            region = (
                "header"
                if row_index == 0
                else "first_column"
                if column_index == 0
                else "body"
            )
            cell_style = styles[region]
            cell.text = labels[row_index][column_index]
            cell.margin_left = Emu(45720)
            cell.margin_right = Emu(45720)
            cell.margin_top = Emu(30000)
            cell.margin_bottom = Emu(30000)
            cell.text_frame.word_wrap = True
            cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            if cell_style["fill_color"] is not None:
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_style["fill_color"]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = cell_style["alignment"]
                for run in paragraph.runs:
                    set_run_fonts(run, cell_style["font_name"])
                    run.font.size = Pt(cell_style["font_size"])
                    run.font.color.rgb = cell_style["text_color"]
                    run.font.bold = cell_style["bold"]


def add_table_from_rows(
    slide: Any,
    box: dict[str, int],
    rows: list[list[str]],
    styles: dict[str, dict[str, Any]],
) -> Any | None:
    usable_rows = [
        [re.sub(r"\s+", " ", cell).strip() for cell in row]
        for row in rows
        if any(cell.strip() for cell in row)
        and not any(
            marker in " ".join(row)
            for marker in ("数据来源", "资料来源")
        )
    ]
    if not usable_rows:
        return None
    column_count = max(len(row) for row in usable_rows)
    if column_count < 2:
        return None
    row_height_pt = (
        box["height"] / EMU_PER_POINT / max(1, len(usable_rows))
    )
    column_width_pt = (
        box["width"] / EMU_PER_POINT / max(1, column_count)
    )
    longest_cell = max(
        (len(cell) for row in usable_rows for cell in row),
        default=1,
    )
    adaptive_font_size = max(
        6.0,
        min(
            14.0,
            row_height_pt * 0.30,
            column_width_pt / max(8.0, longest_cell * 0.50),
        ),
    )
    effective_styles = {
        region: {
            **cell_style,
            "font_size": adaptive_font_size,
        }
        for region, cell_style in styles.items()
    }
    table_shape = slide.shapes.add_table(
        len(usable_rows),
        column_count,
        Emu(box["left"]),
        Emu(box["top"]),
        Emu(box["width"]),
        Emu(box["height"]),
    )
    table_shape.name = "source_table"
    table = table_shape.table
    for row_index, row in enumerate(table.rows):
        source_row = usable_rows[row_index]
        for column_index, cell in enumerate(row.cells):
            region = (
                "header"
                if row_index == 0
                else "first_column"
                if column_index == 0
                else "body"
            )
            cell_style = effective_styles[region]
            cell.text = (
                source_row[column_index]
                if column_index < len(source_row)
                else ""
            )
            cell.margin_left = Emu(45720)
            cell.margin_right = Emu(45720)
            cell.margin_top = Emu(30000)
            cell.margin_bottom = Emu(30000)
            if cell_style["fill_color"] is not None:
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_style["fill_color"]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = cell_style["alignment"]
                for run in paragraph.runs:
                    set_run_fonts(run, cell_style["font_name"])
                    run.font.size = Pt(cell_style["font_size"])
                    run.font.color.rgb = cell_style["text_color"]
                    run.font.bold = cell_style["bold"]
    return table_shape


def remove_template_visual_shapes(slide: Any, configured: Any) -> None:
    if not configured:
        return
    configured_names = configured if isinstance(configured, list) else [configured]
    leaf_names = {name.split("/")[-1] for name in configured_names}
    for top_level in list(slide.shapes):
        descendant_names = {
            shape.name for shape in iter_ppt_shapes([top_level])
        }
        if descendant_names & leaf_names:
            element = top_level.element
            relationship_ids = {
                value
                for descendant in element.iter()
                for attribute, value in descendant.attrib.items()
                if attribute.startswith(
                    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
                )
                and value.startswith("rId")
            }
            element.getparent().remove(element)
            for relationship_id in relationship_ids:
                if relationship_id in slide.part.rels:
                    slide.part.drop_rel(relationship_id)


def visual_seed_box(
    slide: Any,
    configured: Any,
    fallback: dict[str, int],
) -> dict[str, int]:
    if not configured:
        return fallback
    configured_names = configured if isinstance(configured, list) else [configured]
    shapes = [
        find_ppt_shape(slide, name)
        for name in configured_names
    ]
    shapes = [shape for shape in shapes if shape is not None]
    if not shapes:
        return fallback
    left = min(int(shape.left) for shape in shapes)
    top = min(int(shape.top) for shape in shapes)
    right = max(int(shape.left + shape.width) for shape in shapes)
    bottom = max(int(shape.top + shape.height) for shape in shapes)
    return {
        "left": left,
        "top": top,
        "width": right - left,
        "height": bottom - top,
    }


def _chart_asset_partname(relationship_type: str) -> str:
    if relationship_type.endswith("/chartStyle"):
        return "/ppt/charts/style%d.xml"
    if relationship_type.endswith("/chartColorStyle"):
        return "/ppt/charts/colors%d.xml"
    if relationship_type.endswith("/themeOverride"):
        return "/ppt/theme/themeOverride%d.xml"
    if relationship_type.endswith("/chartUserShapes"):
        return "/ppt/drawings/drawing%d.xml"
    return "/ppt/charts/related%d.xml"


def sanitize_chart_extensions(chart_xml: Any) -> Any:
    """Remove producer-private metadata that is not part of the chart visual."""

    chart_namespace = (
        "http://schemas.openxmlformats.org/drawingml/2006/chart"
    )
    wps_namespace = "https://web.wps.cn/et/2018/main"
    for extension in list(
        chart_xml.findall(f".//{{{chart_namespace}}}ext")
    ):
        if any(
            getattr(descendant, "tag", "").startswith(f"{{{wps_namespace}}}")
            for descendant in extension.iterdescendants()
        ):
            parent = extension.getparent()
            if parent is not None:
                parent.remove(extension)
                if (
                    parent.tag == f"{{{chart_namespace}}}extLst"
                    and len(parent) == 0
                ):
                    grandparent = parent.getparent()
                    if grandparent is not None:
                        grandparent.remove(parent)
    return chart_xml


def normalize_chart_axis_ids(
    chart_xml: Any,
    chart_namespace_id: int | None = None,
) -> Any:
    """Convert signed Word chart axis IDs to valid unsigned OOXML values.

    Some Word-authored charts serialize ``c:axId``/``c:crossAx`` as signed
    32-bit integers. PowerPoint expects the same bit pattern as an unsigned
    integer and otherwise reports that the generated presentation is damaged.
    """

    chart_namespace = (
        "http://schemas.openxmlformats.org/drawingml/2006/chart"
    )
    axis_elements = [
        element
        for local_name in ("axId", "crossAx")
        for element in chart_xml.findall(
            f".//{{{chart_namespace}}}{local_name}"
        )
    ]
    for element in axis_elements:
        raw_value = element.get("val")
        if raw_value is None:
            continue
        try:
            numeric_value = int(raw_value)
        except ValueError:
            continue
        if numeric_value < 0:
            element.set("val", str(numeric_value & 0xFFFFFFFF))

    if chart_namespace_id is not None:
        value_map: dict[str, str] = {}
        # Word commonly reuses the same default axis IDs in many charts.
        # PowerPoint treats those cross-chart collisions as damaged content.
        base_value = 2_000_000_000 + chart_namespace_id * 32
        for element in axis_elements:
            raw_value = element.get("val")
            if raw_value is None:
                continue
            if raw_value not in value_map:
                value_map[raw_value] = str(base_value + len(value_map))
            element.set("val", value_map[raw_value])
    return chart_xml


def freeze_source_chart_colors(
    chart_xml: Any,
    chart_item: dict[str, Any],
    project_root: Path,
) -> Any:
    """Resolve source Word theme colors before moving the chart into PPT."""

    c_ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    c = lambda name: f"{{{c_ns}}}{name}"
    a = lambda name: f"{{{a_ns}}}{name}"
    assets = {item.get("kind"): item for item in chart_item.get("related_assets", [])}
    # Some Word/WPS charts do not package a usable themeOverride. In that
    # case PowerPoint otherwise resolves automatic colors against the PPT
    # template (orange/green/magenta). Use the approved GF source palette only
    # as a fallback; explicit RGB series remain untouched below.
    theme_colors: dict[str, str] = {
        "accent1": "7C80C8",
        "accent2": "B8BFE4",
        "accent3": "D9D9D9",
        "accent4": "2E3160",
        "accent5": "AED9FF",
        "accent6": "397099",
    }

    palette_slots: list[str] = []
    colors_asset = assets.get("colors")
    if colors_asset:
        colors_xml = parse_xml(
            (project_root / colors_asset["extracted_path"]).read_bytes()
        )
        palette_slots = [
            node.get("val")
            for node in colors_xml
            if etree.QName(node).localname == "schemeClr"
            and node.get("val") in theme_colors
        ]
    if not palette_slots:
        palette_slots = [
            slot for slot in (
                "accent1", "accent2", "accent3",
                "accent4", "accent5", "accent6",
            )
            if slot in theme_colors
        ]
    if not palette_slots:
        return chart_xml

    def ensure_sppr(series: Any) -> Any:
        sppr = series.find(c("spPr"))
        if sppr is not None:
            return sppr
        sppr = etree.Element(c("spPr"))
        children = list(series)
        insert_at = next(
            (
                index for index, child in enumerate(children)
                if child.tag in {
                    c("marker"), c("dPt"), c("dLbls"),
                    c("cat"), c("val"), c("xVal"), c("yVal"),
                }
            ),
            len(children),
        )
        series.insert(insert_at, sppr)
        return sppr

    def set_solid(parent: Any, color: str) -> None:
        for child in list(parent):
            if child.tag in {
                a("noFill"), a("solidFill"), a("gradFill"), a("pattFill")
            }:
                parent.remove(child)
        solid = etree.Element(a("solidFill"))
        etree.SubElement(solid, a("srgbClr"), val=color)
        parent.insert(0, solid)

    series_index = 0
    plot_area = chart_xml.find(f".//{c('plotArea')}")
    if plot_area is None:
        return chart_xml
    line_types = {"lineChart", "scatterChart", "radarChart", "stockChart"}
    for chart_group in list(plot_area):
        chart_type = etree.QName(chart_group).localname
        for series in chart_group.findall(c("ser")):
            sppr = ensure_sppr(series)
            # Explicit RGB is already source-stable and must remain untouched.
            if sppr.find(f".//{a('srgbClr')}") is not None:
                series_index += 1
                continue
            color = theme_colors[
                palette_slots[series_index % len(palette_slots)]
            ]
            if chart_type in line_types:
                line = sppr.find(a("ln"))
                if line is None:
                    line = etree.SubElement(sppr, a("ln"))
                set_solid(line, color)
            else:
                set_solid(sppr, color)
                line = sppr.find(a("ln"))
                if line is None:
                    line = etree.SubElement(sppr, a("ln"))
                set_solid(line, color)
            series_index += 1
    return chart_xml


def validate_chart_relationships(
    chart_xml: Any,
    chart_part: Any,
) -> None:
    """Reject semantically invalid chart relationships before saving PPTX."""

    relationship_namespace = (
        "http://schemas.openxmlformats.org/"
        "officeDocument/2006/relationships"
    )
    expected_types = {
        "externalData": "/package",
        "userShapes": "/chartUserShapes",
    }
    for element in chart_xml.iter():
        local_name = str(element.tag).rsplit("}", 1)[-1]
        expected_suffix = expected_types.get(local_name)
        if expected_suffix is None:
            continue
        relationship_id = element.get(
            f"{{{relationship_namespace}}}id"
        )
        if not relationship_id or relationship_id not in chart_part.rels:
            raise ValueError(
                f"Chart {local_name} has no valid relationship"
            )
        actual_type = chart_part.rels[relationship_id].reltype
        if not actual_type.endswith(expected_suffix):
            raise ValueError(
                f"Chart {local_name} points to {actual_type}, "
                f"expected *{expected_suffix}"
            )


def set_chart_text_size(chart_xml: Any, font_size_pt: float) -> Any:
    """Set editable chart size plus Latin/East Asian font families."""

    drawing_namespace = (
        "http://schemas.openxmlformats.org/drawingml/2006/main"
    )
    size = str(int(round(font_size_pt * 100)))
    for local_name in ("defRPr", "rPr", "endParaRPr"):
        for element in chart_xml.findall(
            f".//{{{drawing_namespace}}}{local_name}"
        ):
            element.set("sz", size)
            latin = element.find(f"{{{drawing_namespace}}}latin")
            if latin is None:
                latin = etree.SubElement(
                    element,
                    f"{{{drawing_namespace}}}latin",
                )
            latin.set("typeface", "Arial")
            east_asian = element.find(f"{{{drawing_namespace}}}ea")
            if east_asian is None:
                east_asian = etree.Element(f"{{{drawing_namespace}}}ea")
                latin.addnext(east_asian)
            east_asian.set("typeface", "思源黑体 CN Normal")
    return chart_xml


def add_docx_chart(
    slide: Any,
    chart_item: dict[str, Any],
    box: dict[str, int],
    project_root: Path,
    font_size_pt: float,
) -> Any:
    """Transfer a Word chart as a native PowerPoint chart while filling layout box."""

    chart_path = project_root / chart_item["extracted_path"]
    workbook_path = project_root / chart_item["workbook_path"]
    chart_xml = parse_xml(chart_path.read_bytes())
    sanitize_chart_extensions(chart_xml)
    normalize_chart_axis_ids(chart_xml)
    freeze_source_chart_colors(chart_xml, chart_item, project_root)
    set_chart_text_size(chart_xml, font_size_pt)

    chart_data = CategoryChartData()
    chart_data.categories = ["原文数据"]
    chart_data.add_series("原文序列", (0.0,))

    # ============================================================
    # 核心修改：
    # 不再根据 Word 原始 chart 尺寸缩放
    # 直接使用 PPT layout engine 分配出的 box
    # 避免生成后图表缩成左上角小块
    # ============================================================

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Emu(box["left"]),
        Emu(box["top"]),
        Emu(box["width"]),
        Emu(box["height"]),
        chart_data,
    )

    chart_shape.name = "source_docx_chart"

    chart_part = chart_shape.chart.part
    chart_part_number_match = re.search(
        r"(\d+)\.xml$",
        str(chart_part.partname),
    )
    if chart_part_number_match:
        normalize_chart_axis_ids(
            chart_xml,
            int(chart_part_number_match.group(1)),
        )


    package_relationship = next(
        (
            (relationship_id, relationship)
            for relationship_id, relationship in chart_part.rels.items()
            if relationship.reltype.endswith("/package")
        ),
        None,
    )

    if package_relationship is None:
        raise ValueError(
            "Temporary PowerPoint chart has no embedded workbook"
        )


    workbook_relationship_id, workbook_relationship = package_relationship

    workbook_relationship.target_part._blob = (
        workbook_path.read_bytes()
    )


    relationship_namespace = (
        "http://schemas.openxmlformats.org/"
        "officeDocument/2006/relationships"
    )


    package = chart_part.package


    relationship_id_map: dict[str, str] = {}
    for asset in chart_item.get("related_assets", []):

        asset_path = (
            project_root /
            asset["extracted_path"]
        )

        asset_part = Part(
            package.next_partname(
                _chart_asset_partname(
                    asset["relationship_type"]
                )
            ),
            asset["content_type"],
            package,
            asset_path.read_bytes(),
        )

        new_relationship_id = chart_part.relate_to(
            asset_part,
            asset["relationship_type"],
        )
        source_relationship_id = asset.get("relationship_id")
        if source_relationship_id:
            relationship_id_map[str(source_relationship_id)] = (
                new_relationship_id
            )

    # The chart XML came from Word, while the destination relationships were
    # allocated by python-pptx.  Their rIds are independent.  Leaving the Word
    # IDs in place can make (for example) c:userShapes point at a themeOverride
    # part, which PowerPoint treats as damaged content and removes during its
    # repair pass.
    for descendant in chart_xml.iter():
        for attribute, value in list(descendant.attrib.items()):
            if (
                attribute.startswith(f"{{{relationship_namespace}}}")
                and value in relationship_id_map
            ):
                descendant.set(
                    attribute,
                    relationship_id_map[value],
                )

    # Set the workbook relationship last. Destination rIds can coincidentally
    # equal a source asset rId, so setting it earlier could redirect
    # externalData to a chart style or user-shapes part during the rewrite.
    for external_data in chart_xml.findall(
        ".//{http://schemas.openxmlformats.org/"
        "drawingml/2006/chart}externalData"
    ):
        external_data.set(
            f"{{{relationship_namespace}}}id",
            workbook_relationship_id,
        )

    validate_chart_relationships(chart_xml, chart_part)

    chart_part._element = chart_xml

    chart_part.__dict__.pop(
        "chart",
        None,
    )

    chart_part.__dict__.pop(
        "chart_workbook",
        None,
    )


    return chart_shape


def split_visual_boxes(
    box: dict[str, int],
    count: int,
) -> list[dict[str, int]]:
    if count <= 1:
        return [box]
    gap = 100000
    if count == 2:
        width = (box["width"] - gap) // 2
        return [
            {
                "left": box["left"] + index * (width + gap),
                "top": box["top"],
                "width": width,
                "height": box["height"],
            }
            for index in range(2)
        ]
    rows = math.ceil(count / 2)
    width = (box["width"] - gap) // 2
    height = (box["height"] - gap * (rows - 1)) // rows
    return [
        {
            "left": box["left"] + (index % 2) * (width + gap),
            "top": box["top"] + (index // 2) * (height + gap),
            "width": width,
            "height": height,
        }
        for index in range(count)
    ]


def grouped_table_boxes(
    box: dict[str, int],
    visual_items: list[tuple[str, Any, str]],
) -> list[dict[str, int]] | None:
    if len(visual_items) < 2 or any(
        kind != "table" for kind, _, _ in visual_items
    ):
        return None
    group_indexes = {
        item.get("wrapper_group_index") for _, item, _ in visual_items
    }
    if len(group_indexes) != 1 or None in group_indexes:
        return None
    gap = 100000
    available_height = box["height"] - gap * (len(visual_items) - 1)
    row_counts = [
        max(1, int(item.get("row_count", 1)))
        for _, item, _ in visual_items
    ]
    total_rows = sum(row_counts)
    boxes: list[dict[str, int]] = []
    top = box["top"]
    remaining_height = available_height
    for index, row_count in enumerate(row_counts):
        height = (
            remaining_height
            if index == len(row_counts) - 1
            else available_height * row_count // total_rows
        )
        boxes.append(
            {
                "left": box["left"],
                "top": top,
                "width": box["width"],
                "height": height,
            }
        )
        top += height + gap
        remaining_height -= height
    return boxes


def add_fitted_picture(
    slide: Any,
    path: Path,
    box: dict[str, int],
    width_emu: int | None = None,
    height_emu: int | None = None,
    fill_box: bool = False,
) -> Any:
    file_size = path.stat().st_size
    if file_size > MAX_EMBEDDED_IMAGE_BYTES:
        raise EmbeddedImageError(
            f"图片“{path.name}”文件过大（{file_size / 1024 / 1024:.1f} MB），"
            "请在Word中压缩该图片后重新上传。"
        )

    # Pillow checks the declared pixel count while opening the header, even
    # though neither this function nor python-pptx needs to decode the raster.
    # Temporarily lift that global check under a lock, validate against our own
    # hard limits, and embed the original bytes without resampling.
    with allow_large_image_metadata():
        try:
            with Image.open(path) as image:
                pixel_width, pixel_height = image.size
        except (OSError, ValueError) as exc:
            raise EmbeddedImageError(
                f"图片“{path.name}”无法识别或已损坏，请替换后重新上传。"
            ) from exc

        pixel_count = pixel_width * pixel_height
        if pixel_count > MAX_EMBEDDED_IMAGE_PIXELS:
            raise EmbeddedImageError(
                f"图片“{path.name}”分辨率过高"
                f"（{pixel_width}×{pixel_height}，约{pixel_count / 1_000_000:.0f}百万像素），"
                "请在Word中压缩图片至220ppi或更低后重新上传。"
            )

        # The Word drawing extent is the authoritative display aspect ratio.
        # Pixel dimensions are only a fallback for malformed/missing extents.
        source_width = int(width_emu or pixel_width)
        source_height = int(height_emu or pixel_height)
        if fill_box:
            # A single raster visual follows the sell-side template's
            # full-width treatment: fill the allocated visual region exactly,
            # aligned with the body bullet on the left and the same margin on
            # the right. Multi-image pages continue to preserve aspect ratio.
            width = box["width"]
            height = box["height"]
        else:
            scale = min(
                box["width"] / max(1, source_width),
                box["height"] / max(1, source_height),
            )
            width = int(source_width * scale)
            height = int(source_height * scale)
        shape = slide.shapes.add_picture(
            str(path),
            Emu(
                box["left"]
                if fill_box
                else box["left"] + (box["width"] - width) // 2
            ),
            Emu(box["top"]),
            Emu(width),
            Emu(height),
        )
    shape.name = "source_docx_image"
    return shape


def visual_caption(
    mapping: dict[str, Any],
    table: dict[str, Any] | None,
    kind: str,
    source_id: int | None,
) -> str:
    """Return the source caption, omitting the report's figure number."""

    captions: list[str] = []
    if table is not None:
        texts = [str(table.get("caption") or "")]
        texts.extend(
            str(cell)
            for row in table.get("wrapper_rows", table.get("rows", []))
            for cell in row
        )
        captions = [
            re.sub(r"\s+", " ", text).strip()
            for text in texts
            if re.match(r"^\s*(?:图|表)\s*\d+\s*[：:]", str(text))
        ]
        captions = list(dict.fromkeys(captions))
    if captions:
        source_ids = (
            table.get("chart_indexes", [])
            if kind == "chart"
            else table.get("image_indexes", [])
            if kind == "image"
            else []
        )
        position = source_ids.index(source_id) if source_id in source_ids else 0
        caption = captions[min(position, len(captions) - 1)]
    else:
        candidate = str(mapping.get("paragraph_text") or "").strip()
        if (
            mapping.get("mapping_basis") != "visual_title"
            and not re.match(r"^(?:图|表)\s*\d+\s*[：:]", candidate)
        ):
            return ""
        caption = candidate
    # A Word wrapper cell can contain both the caption and a following source
    # line.  Only the caption belongs in the grey title strip.
    caption = re.split(
        r"\s*(?:数据|资料)来源\s*[：:]?",
        caption,
        maxsplit=1,
    )[0].strip()
    if not caption:
        return ""
    label = "表" if kind == "table" else "图"
    caption = re.sub(
        r"^(?:图|表)\s*\d+\s*[：:]\s*",
        f"{label} ：",
        caption,
    )
    if not re.match(r"^(?:图|表)\s*[：:]", caption):
        caption = f"{label} ：{caption}"
    return caption.rstrip("。")


def add_visual_title(
    slide: Any,
    text: str,
    box: dict[str, int],
) -> Any | None:
    if not text:
        return None
    shape = slide.shapes.add_textbox(
        Emu(box["left"]),
        Emu(box["top"]),
        Emu(box["width"]),
        Emu(box["height"]),
    )
    shape.name = "source_visual_title"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string("F2F2F2")
    shape.line.fill.background()
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    set_run_fonts(run, "Source Han Sans CN Regular")
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("000000")
    return shape


def add_source_visuals(
    slide: Any,
    visual_mappings: list[dict[str, Any]],
    document: dict[str, Any],
    box: dict[str, int],
    table_styles: dict[str, dict[str, Any]],
    project_root: Path,
    single_visual_box: dict[str, int] | None = None,
) -> int:
    tables = {item["index"]: item for item in document.get("tables", [])}
    images = {item["index"]: item for item in document.get("images", [])}
    charts = {item["index"]: item for item in document.get("charts", [])}
    visual_items: list[tuple[str, Any, str]] = []
    for mapping in visual_mappings:
        table = tables.get(mapping.get("table_id"))
        chart_ids = list(mapping.get("embedded_chart_ids", []))
        image_ids = list(mapping.get("embedded_image_ids", []))
        if table is not None:
            chart_ids = chart_ids or table.get("chart_indexes", [])
            image_ids = image_ids or table.get("image_indexes", [])
        if chart_ids:
            visual_items.extend(
                (
                    "chart",
                    charts[index],
                    visual_caption(mapping, table, "chart", index),
                )
                for index in chart_ids
                if index in charts
            )
        elif image_ids:
            visual_items.extend(
                (
                    "image",
                    images[index],
                    visual_caption(mapping, table, "image", index),
                )
                for index in image_ids
                if index in images
            )
        elif mapping.get("image_id") in images:
            image_id = mapping["image_id"]
            visual_items.append(
                (
                    "image",
                    images[image_id],
                    visual_caption(mapping, table, "image", image_id),
                )
            )
        elif table is not None:
            visual_items.append(
                (
                    "table",
                    table,
                    visual_caption(mapping, table, "table", table["index"]),
                )
            )
    layout_box = dict(box)
    grouped_tables = grouped_table_boxes(layout_box, visual_items)
    is_single_visual_group = len(visual_items) == 1 or grouped_tables is not None
    if is_single_visual_group and single_visual_box is not None:
        layout_box.update(single_visual_box)
        grouped_tables = grouped_table_boxes(layout_box, visual_items)
    boxes = grouped_tables or split_visual_boxes(
        layout_box, len(visual_items)
    )
    rendered = 0
    chart_font_size = 12 if len(visual_items) == 1 else 10
    title_height = 261610
    title_gap = 45000
    for (kind, item, caption), item_box in zip(visual_items, boxes):
        content_box = dict(item_box)
        if caption:
            content_box["top"] += title_height + title_gap
            content_box["height"] = max(
                1,
                content_box["height"] - title_height - title_gap,
            )
        visual_shape = None
        if kind == "chart":
            visual_shape = add_docx_chart(
                slide,
                item,
                content_box,
                project_root,
                chart_font_size,
            )
            rendered += 1
        elif kind == "image":
            visual_shape = add_fitted_picture(
                slide,
                project_root / item["extracted_path"],
                content_box,
                item.get("width_emu"),
                item.get("height_emu"),
                fill_box=is_single_visual_group,
            )
            rendered += 1
        else:
            visual_shape = add_table_from_rows(
                slide,
                content_box,
                item.get("rows", []),
                table_styles,
            )
            if visual_shape is not None:
                rendered += 1
        if caption and visual_shape is not None:
            add_visual_title(
                slide,
                caption,
                {
                    "left": (
                        int(visual_shape.left)
                        if kind == "image"
                        else item_box["left"]
                    ),
                    "top": item_box["top"],
                    "width": (
                        int(visual_shape.width)
                        if kind == "image"
                        else item_box["width"]
                    ),
                    "height": title_height,
                },
            )
    return rendered


def choose_template_examples(mapping: dict[str, Any]) -> dict[str, list[dict[str, Any]]]:
    grouped: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for item in mapping["slides"]:
        grouped[item["page_type"]].append(item)
    return grouped


def resolve_template_path(data: dict[str, Any]) -> Path:
    configured = data.get("template_path")
    if configured:
        path = Path(configured)
    else:
        source_template = data.get("style", {}).get("source_template")
        if not source_template:
            raise ValueError("style_config.yaml does not declare source_template")
        path = Path(__file__).resolve().parent.parent / "template" / source_template
    path = path.resolve()
    if not path.is_file():
        raise FileNotFoundError(f"Template PPTX not found: {path}")
    return path


def select_template_mapping(
    page_type: str,
    examples: dict[str, list[dict[str, Any]]],
    example_usage: dict[str, int],
) -> dict[str, Any]:
    candidates = examples.get(page_type)
    if not candidates:
        raise ValueError(f"No template mapping available for page type: {page_type}")
    # These semantic page types use one stable example selected by
    # slide_mapping.yaml. Their template position may change independently of
    # the generated presentation order, so never hard-code a template slide
    # number here. For section, using the first mapped example preserves the
    # existing fixed section design without rotating across section variants.
    if page_type in {"cover", "summary", "section", "risk", "thanks"}:
        return candidates[0]
    candidate_index = example_usage[page_type] % len(candidates)
    mapped = candidates[candidate_index]
    example_usage[page_type] += 1
    return mapped


def order_generation_plans(
    plan_slides: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    """Sort planned slides normally while always placing thanks pages last."""
    return sorted(
        plan_slides,
        key=lambda plan: (
            plan.get("page_type") == "thanks",
            plan["slide_number"],
        ),
    )


def remove_initial_slides(presentation: Presentation, count: int) -> None:
    for _ in range(count):
        slide_id = presentation.slides._sldIdLst[0]
        presentation.part.drop_rel(slide_id.rId)
        del presentation.slides._sldIdLst[0]


def clone_text_shape(
    source_shape: Any,
    target_slide: Any,
    text: str,
    *,
    font_name: str,
    font_size: float,
    color: RGBColor,
    bold: bool,
    height: int | None = None,
) -> Any:
    element = deepcopy(source_shape.element)
    target_slide.shapes._spTree.insert_element_before(element, "p:extLst")
    shape = target_slide.shapes[-1]
    if height is not None:
        shape.height = Emu(height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    for run in paragraph.runs:
        set_run_fonts(run, font_name)
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
    return shape


def iter_ppt_shapes(shapes: Any):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_ppt_shapes(shape.shapes)


def remove_all_slide_shapes(slide: Any) -> None:
    for shape in list(slide.shapes):
        element = shape.element
        element.getparent().remove(element)


def copy_slide_background(source_slide: Any, target_slide: Any) -> None:
    source_background = source_slide._element.cSld.bg
    if source_background is None:
        return
    target_background = target_slide._element.cSld.bg
    if target_background is not None:
        target_slide._element.cSld.remove(target_background)
    target_slide._element.cSld.insert(0, deepcopy(source_background))


def duplicate_template_slide(
    presentation: Presentation,
    source_slide: Any,
) -> Any:
    """Clone a complete template example slide inside the same presentation."""
    target_slide = presentation.slides.add_slide(source_slide.slide_layout)
    # add_slide() creates layout placeholders. Remove them before copying the
    # source slide so no empty "Click to add title" placeholders survive.
    remove_all_slide_shapes(target_slide)

    relationship_map: dict[str, str] = {}
    for relationship in source_slide.part.rels.values():
        if relationship.reltype.endswith("/slideLayout") or relationship.reltype.endswith(
            "/notesSlide"
        ):
            continue
        relationship_map[relationship.rId] = target_slide.part.relate_to(
            relationship.target_ref if relationship.is_external else relationship.target_part,
            relationship.reltype,
            relationship.is_external,
        )

    for source_shape in source_slide.shapes:
        element = deepcopy(source_shape.element)
        for descendant in element.iter():
            for attribute, value in list(descendant.attrib.items()):
                if value in relationship_map:
                    descendant.set(attribute, relationship_map[value])
            if descendant.tag.endswith("}creationId"):
                if "id" in descendant.attrib:
                    descendant.set("id", "{" + str(uuid.uuid4()).upper() + "}")
                if "val" in descendant.attrib:
                    descendant.set("val", str(secrets.randbits(32)))
            elif descendant.tag.endswith("}modId") and "val" in descendant.attrib:
                descendant.set("val", str(secrets.randbits(32)))
        target_slide.shapes._spTree.insert_element_before(element, "p:extLst")

    copy_slide_background(source_slide, target_slide)
    return target_slide


def find_ppt_shape(slide: Any, configured: Any) -> Any | None:
    if not configured:
        return None
    names = configured if isinstance(configured, list) else [configured]
    leaf_names = {name.split("/")[-1] for name in names}
    return next(
        (
            shape
            for shape in iter_ppt_shapes(slide.shapes)
            if shape.name in leaf_names
        ),
        None,
    )


def find_source_shape(slide: Any) -> Any | None:
    return next(
        (
            shape
            for shape in iter_ppt_shapes(slide.shapes)
            if getattr(shape, "has_text_frame", False)
            and (shape.text or "").strip().startswith("数据来源")
        ),
        None,
    )


def find_page_number_shape(slide: Any) -> Any | None:
    return next(
        (
            shape
            for shape in iter_ppt_shapes(slide.shapes)
            if "编号" in shape.name
        ),
        None,
    )


def find_text_shape_matching(slide: Any, pattern: str) -> Any | None:
    regex = re.compile(pattern)
    return next(
        (
            shape
            for shape in iter_ppt_shapes(slide.shapes)
            if getattr(shape, "has_text_frame", False)
            and regex.search((shape.text or "").strip())
        ),
        None,
    )


def text_style_snapshot(shape: Any) -> dict[str, Any]:
    if not getattr(shape, "has_text_frame", False):
        return {}
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            try:
                font_color = run.font.color.rgb
            except (AttributeError, TypeError):
                font_color = None
            return {
                "font_name": run.font.name,
                "font_size": run.font.size,
                "font_bold": run.font.bold,
                "font_color": font_color,
                "alignment": paragraph.alignment,
            }
    return {}


def replace_shape_text(
    shape: Any | None,
    text: str,
    *,
    bullet: bool = False,
    font_name: str | None = None,
    font_size: float | None = None,
    color: RGBColor | None = None,
    bold: bool | None = None,
) -> None:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    saved = text_style_snapshot(shape)
    frame = shape.text_frame
    bullet_properties = (
        deepcopy(frame.paragraphs[0]._p.get_or_add_pPr()) if bullet else None
    )
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    lines = text.splitlines() if text else [""]
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        # The template body placeholder already supplies the square bullet and
        # indentation. Copy that paragraph formatting to every independent
        # point; a literal "•" would create a double bullet on the first one.
        if bullet_properties is not None:
            if paragraph._p.pPr is not None:
                paragraph._p.remove(paragraph._p.pPr)
            paragraph._p.insert(0, deepcopy(bullet_properties))
        if saved.get("alignment") is not None:
            paragraph.alignment = saved["alignment"]
        effective_name = font_name or saved.get("font_name")
        effective_size = Pt(font_size) if font_size is not None else saved.get("font_size")
        effective_color = color or saved.get("font_color")
        effective_bold = bold if bold is not None else saved.get("font_bold")
        for segment in script_segments(line):
            run = paragraph.add_run()
            run.text = segment
            if effective_name:
                set_run_fonts(run, effective_name)
            if effective_size:
                run.font.size = effective_size
            if effective_color:
                run.font.color.rgb = effective_color
            if effective_bold is not None:
                run.font.bold = effective_bold


def align_bullet_text_with_title(
    content_shape: Any | None,
    title_shape: Any | None,
) -> None:
    """Align bullet text with the visible title start, keeping bullets left."""

    if (
        content_shape is None
        or title_shape is None
        or not getattr(content_shape, "has_text_frame", False)
        or not getattr(title_shape, "has_text_frame", False)
    ):
        return
    content_frame = content_shape.text_frame
    title_frame = title_shape.text_frame
    title_text_left = (
        int(title_shape.left) + int(title_frame.margin_left or 0)
    )
    content_inner_left = (
        int(content_shape.left) + int(content_frame.margin_left or 0)
    )
    text_indent = max(0, title_text_left - content_inner_left)
    for paragraph in content_frame.paragraphs:
        properties = paragraph._p.get_or_add_pPr()
        properties.set("marL", str(text_indent))
        properties.set("indent", str(-text_indent))


def bullet_left_edge(content_shape: Any | None) -> int | None:
    """Return the visible bullet-square left edge for visual alignment."""

    if content_shape is None or not getattr(
        content_shape, "has_text_frame", False
    ):
        return None
    return int(content_shape.left) + int(
        content_shape.text_frame.margin_left or 0
    )


def split_summary_point(text: str) -> tuple[str, str]:
    """Split a conclusion into a red lead and black explanatory body."""
    normalized = re.sub(r"\s+", " ", text).strip()
    colon_match = re.match(r"^(.{2,24}?[：:])\s*(.+)$", normalized)
    if colon_match:
        return colon_match.group(1), colon_match.group(2)
    for separator in ("，", "；", "。"):
        position = normalized.find(separator)
        if 6 <= position <= 22 and position < len(normalized) - 1:
            return normalized[: position + 1], normalized[position + 1 :].lstrip()
    split_at = min(14, len(normalized))
    return normalized[:split_at], normalized[split_at:]


def unique_points(points: list[str]) -> list[str]:
    """Remove repeated bullets while preserving source order."""
    seen: set[str] = set()
    result: list[str] = []
    for point in points:
        normalized = re.sub(r"^[\s•·▪■]+", "", str(point)).strip()
        comparison_key = re.sub(r"\s+", "", normalized).rstrip("。；;")
        if not normalized or comparison_key in seen:
            continue
        seen.add(comparison_key)
        result.append(normalized)
    return result


def replace_summary_text(
    shape: Any | None,
    points: list[str],
    rich_text: list[dict[str, Any]] | None = None,
) -> None:
    """Populate summary text, coloring only emphasis copied from the Word source."""
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    frame = shape.text_frame
    frame.word_wrap = True
    frame.auto_size = None
    paragraphs = frame.paragraphs
    bullet_properties = deepcopy(paragraphs[0]._p.get_or_add_pPr())
    spacer_properties = (
        deepcopy(paragraphs[1]._p.get_or_add_pPr())
        if len(paragraphs) > 1 and not paragraphs[1].text.strip()
        else None
    )
    for paragraph in paragraphs:
        paragraph.clear()
    source_paragraphs = rich_text or [
        {
            "text": point,
            "runs": [{"text": point, "bold": False}],
        }
        for point in points
    ]
    for point_index, source_paragraph in enumerate(source_paragraphs):
        # The normative summary layout alternates a bullet paragraph and one
        # genuinely blank paragraph.  Recreate that rhythm instead of relying
        # on a small ``space_after`` value, which is rendered too tightly by
        # PowerPoint when a bullet wraps onto multiple lines.
        paragraph_index = point_index * 2
        while paragraph_index >= len(frame.paragraphs):
            frame.add_paragraph()
        paragraph = frame.paragraphs[paragraph_index]
        if paragraph._p.pPr is not None:
            paragraph._p.remove(paragraph._p.pPr)
        paragraph._p.insert(0, deepcopy(bullet_properties))
        paragraph.line_spacing = 1.5
        paragraph.space_after = Pt(0)
        runs = source_paragraph.get("runs") or [
            {"text": source_paragraph.get("text", ""), "bold": False}
        ]
        for source_run in runs:
            text = str(source_run.get("text", ""))
            if not text:
                continue
            emphasized = bool(source_run.get("bold"))
            target_run = paragraph.add_run()
            target_run.text = text
            set_run_fonts(
                target_run,
                "思源黑体 CN Bold"
                if emphasized
                else "思源黑体 CN Regular",
            )
            target_run.font.size = Pt(12)
            target_run.font.bold = emphasized
            target_run.font.color.rgb = RGBColor.from_string(
                "C00000" if emphasized else "000000"
            )
        if point_index < len(source_paragraphs) - 1:
            spacer_index = paragraph_index + 1
            while spacer_index >= len(frame.paragraphs):
                frame.add_paragraph()
            spacer = frame.paragraphs[spacer_index]
            spacer.clear()
            if spacer._p.pPr is not None:
                spacer._p.remove(spacer._p.pPr)
            if spacer_properties is not None:
                spacer._p.insert(0, deepcopy(spacer_properties))
            spacer.line_spacing = 1.5
            spacer.space_before = Pt(0)
            spacer.space_after = Pt(0)


def _replace_paragraph_text(
    paragraph: Any,
    text: str,
    *,
    font_name: str | None = None,
    font_size: float | None = None,
    color: RGBColor | None = None,
    bold: bool | None = None,
) -> None:
    paragraph.clear()
    if not text:
        return
    run = paragraph.add_run()
    run.text = text
    if font_name:
        set_run_fonts(run, font_name)
    if font_size is not None:
        run.font.size = Pt(font_size)
    if color is not None:
        run.font.color.rgb = color
    if bold is not None:
        run.font.bold = bold


def replace_risk_content(
    shape: Any | None,
    risk_points: list[str],
    _disclaimer: str,
) -> None:
    """Fill the inherited risk template without flattening its hierarchy."""
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    frame = shape.text_frame
    while len(frame.paragraphs) < 8:
        frame.add_paragraph()
    risk_text = "；".join(
        point.strip().rstrip("。；;") for point in risk_points if point.strip()
    )
    if risk_text and not risk_text.endswith("。"):
        risk_text += "。"
    _replace_paragraph_text(
        frame.paragraphs[0],
        "风险提示",
        font_name="思源黑体 CN Normal",
        color=RGBColor.from_string("666666"),
        bold=False,
    )
    _replace_paragraph_text(
        frame.paragraphs[1],
        risk_text or "风险提示内容以原始研究报告为准。",
        font_name="思源黑体 CN Normal",
        font_size=10,
        color=RGBColor.from_string("666666"),
        bold=False,
    )
    frame.paragraphs[2].clear()
    _replace_paragraph_text(
        frame.paragraphs[3],
        "免责声明",
        font_name="思源黑体 CN Normal",
        color=RGBColor.from_string("666666"),
        bold=False,
    )
    for offset, paragraph in enumerate(frame.paragraphs[4:]):
        text = (
            FIXED_DISCLAIMER_PARAGRAPHS[offset]
            if offset < len(FIXED_DISCLAIMER_PARAGRAPHS)
            else ""
        )
        _replace_paragraph_text(
            paragraph,
            text,
            font_name="Source Han Sans CN Normal",
            font_size=8,
            color=RGBColor.from_string("666666"),
            bold=False,
        )


def clear_unassigned_template_text(
    slide: Any,
    keep_shape_ids: set[int],
) -> None:
    for shape in iter_ppt_shapes(slide.shapes):
        if id(shape.element) in keep_shape_ids:
            continue
        if getattr(shape, "has_text_frame", False):
            replace_shape_text(shape, "")


def template_source_text(content: dict[str, Any]) -> str:
    """Return only the audience-facing source line.

    ``source`` and ``source_trace`` are internal provenance fields and must
    never expose DOCX filenames or heading paths on a client-facing slide.
    """
    return str(content.get("display_source") or "").strip()


def build_presentation(
    data: dict[str, Any],
    layout_debug: list[dict[str, Any]] | None = None,
) -> Presentation:
    validate_inputs(data)
    plan_slides = data["plan"]["slides"]
    content_by_number = {
        item["slide_number"]: item for item in data["content"]["slides"]
    }
    style = data["style"]
    style_values = style["style"]
    layout_file = data["layout"]["files"][0]
    layouts_by_number = {
        item["slide_number"]: item for item in layout_file["slides"]
    }
    examples = choose_template_examples(data["mapping"])
    example_usage: dict[str, int] = defaultdict(int)

    template_path = resolve_template_path(data)
    prs = Presentation(str(template_path))
    template_slide_count = len(prs.slides)
    template_slides = list(prs.slides)

    fonts = style_values["font"]
    colors = style_values["color"]
    sizes = style_values["size"]
    title_color = rgb(colors["title_color"])
    body_color = rgb(colors["body_color"])
    table_text_color = RGBColor.from_string("000000")
    table_fill_color = RGBColor.from_string("2E3160")
    table_font = "思源黑体 CN Medium"
    table_size = numeric(sizes["table_text_pt"], 7)
    table_styles = {
        "header": table_cell_style(
            style_values,
            "header",
            fallback_font=table_font,
            fallback_size=table_size,
            fallback_text_color=RGBColor(255, 255, 255),
            fallback_fill_color=table_fill_color,
            fallback_bold=True,
            fallback_alignment=PP_ALIGN.CENTER,
        ),
        "first_column": table_cell_style(
            {},
            "first_column",
            fallback_font=table_font,
            fallback_size=table_size,
            fallback_text_color=table_text_color,
            fallback_fill_color=RGBColor.from_string("F2F2F2"),
            fallback_bold=False,
            fallback_alignment=PP_ALIGN.CENTER,
        ),
        "body": table_cell_style(
            {},
            "body",
            fallback_font=table_font,
            fallback_size=table_size,
            fallback_text_color=table_text_color,
            fallback_fill_color=RGBColor.from_string("FFFFFF"),
            fallback_bold=False,
            fallback_alignment=PP_ALIGN.CENTER,
        ),
    }
    table_styles["header"] = table_cell_style(
        {},
        "header",
        fallback_font=table_font,
        fallback_size=table_size,
        fallback_text_color=RGBColor(255, 255, 255),
        fallback_fill_color=table_fill_color,
        fallback_bold=False,
        fallback_alignment=PP_ALIGN.CENTER,
    )

    section_number = 0
    document = (
        data.get("doc", {}).get("files", [{}])[0]
        if data.get("doc", {}).get("files")
        else {}
    )
    report_metadata = document.get("report_metadata", {})

    ordered_plans = order_generation_plans(plan_slides)

    for plan in ordered_plans:   
        page_type = plan["page_type"]
        mapped = select_template_mapping(page_type, examples, example_usage)
        source_template_slide = template_slides[mapped["slide_number"] - 1]
        content = content_by_number[plan["slide_number"]]
        slide = duplicate_template_slide(prs, source_template_slide)
        title_shape = find_ppt_shape(slide, mapped.get("title_shape"))
        content_shape = find_ppt_shape(slide, mapped.get("content_shape"))
        source_shape = find_source_shape(slide)
        page_number_shape = find_page_number_shape(slide)
        configured_visuals = [
            *(
                mapped.get("chart_shape")
                if isinstance(mapped.get("chart_shape"), list)
                else [mapped.get("chart_shape")]
                if mapped.get("chart_shape")
                else []
            ),
            *(
                mapped.get("table_shape")
                if isinstance(mapped.get("table_shape"), list)
                else [mapped.get("table_shape")]
                if mapped.get("table_shape")
                else []
            ),
        ]
        visual_box = visual_seed_box(
            slide,
            configured_visuals,
            {
                "left": 650000,
                "top": 2600000,
                "width": int(prs.slide_width) - 1300000,
                "height": int(prs.slide_height) - 3200000,
            },
        )
        title_font, title_size = title_style_for_page(page_type, fonts, sizes)
        points = unique_points(content.get("key_points", []))
        body_size = numeric(sizes["body_text_pt"], 14)
        rendered_visual_count = 0

        if page_type == "cover":
            # Keep template subtitle, author, date, imagery and decoration. Only
            # replace the cover title; planned key_points are intentionally ignored.
            replace_shape_text(
                title_shape,
                content["title"],
                font_name=title_font,
                font_size=title_size,
                color=title_color,
                bold=True,
            )
            report_date = report_metadata.get("report_date")
            if report_date:
                date_match = re.match(r"(\d{4})[/-](\d{1,2})", report_date)
                date_text = (
                    f"{date_match.group(1)}年{int(date_match.group(2))}月"
                    if date_match
                    else report_date
                )
                replace_shape_text(
                    find_text_shape_matching(slide, r"^\d{4}年\d{1,2}月$"),
                    date_text,
                )
            authors = report_metadata.get("authors", [])
            author_shape = find_text_shape_matching(slide, r"分析师|执证号")
            if authors and author_shape is not None:
                existing_lines = [
                    line.strip()
                    for line in (author_shape.text or "").splitlines()
                    if line.strip()
                ]
                retained = [
                    line
                    for line in existing_lines
                    if any(line.startswith(author) for author in authors)
                ]
                replace_shape_text(
                    author_shape,
                    "\n".join(retained or authors),
                )
            replace_shape_text(source_shape, "")
            replace_shape_text(page_number_shape, "")
        elif page_type == "summary":
            clear_unassigned_template_text(
                slide,
                {
                    id(shape.element)
                    for shape in (title_shape, content_shape, page_number_shape)
                    if shape is not None
                },
            )
            replace_shape_text(
                title_shape,
                content["title"],
                font_name=title_font,
                font_size=title_size,
                color=title_color,
                bold=True,
            )
            replace_summary_text(
                content_shape,
                points,
                content.get("rich_text") or None,
            )

            replace_shape_text(page_number_shape, str(plan["slide_number"]))
            if content_shape is not None and page_number_shape is not None:
                maximum_height = (
                    int(page_number_shape.top) - int(content_shape.top) - 90000
                )
                if maximum_height > 0:
                    content_shape.height = Emu(
                        min(int(content_shape.height), maximum_height)
                    )
        elif page_type == "section":
            section_number += 1
            clear_unassigned_template_text(
                slide,
                {
                    id(shape.element)
                    for shape in (title_shape, content_shape)
                    if shape is not None
                },
            )
            section_display_title = content["title"]
            if content.get("subtitle"):
                section_display_title += f"\n——{content['subtitle']}"
            replace_shape_text(
                title_shape,
                section_display_title,
                font_name=fonts["section_title_style"],
                font_size=numeric(sizes["section_title_pt"], 40),
                color=title_color,
                bold=True,
            )
            if title_shape is not None:
                # The template uses a lower baseline for one-line section
                # titles (slide 23) and a higher box for two-line titles
                # (slide 3). Preserve that distinction while retaining the
                # fixed section background/decorations from slide 3.
                if len(re.sub(r"\s+", "", section_display_title)) <= 20:
                    title_shape.top = Emu(1736726)
                    title_shape.height = Emu(799645)
                else:
                    title_shape.top = Emu(1120000)
                    title_shape.height = Emu(1450000)
            replace_shape_text(content_shape, f"{section_number:02d}")
            replace_shape_text(source_shape, "")
            replace_shape_text(page_number_shape, "")
        elif page_type == "risk":
            clear_unassigned_template_text(
                slide,
                {
                    id(shape.element)
                    for shape in (title_shape, content_shape)
                    if shape is not None
                },
            )
            replace_shape_text(
                title_shape,
                "风险提示及免责声明",
            )
            replace_risk_content(
                content_shape,
                points,
                report_metadata.get("legal_disclaimer") or "",
            )
            replace_shape_text(source_shape, "")
            replace_shape_text(page_number_shape, "")
        
        elif page_type == "thanks":
            # The closing page is entirely defined by its inherited template
            # layout ("结尾页"). Keep it byte-for-byte at slide level: even
            # apparently empty local shapes are intentional off-canvas palette
            # helpers, while the visible Thanks artwork lives on the layout.
            pass
        
        else:
            allowed_source = page_type in {
                "chart_analysis",
                "chart_table",
                "table_summary",
                "matrix",
            }
            keep = {
                id(shape.element)
                for shape in (
                    title_shape,
                    content_shape,
                    source_shape if allowed_source else None,
                    page_number_shape,
                )
                if shape is not None
            }
            clear_unassigned_template_text(slide, keep)
            remove_template_visual_shapes(slide, configured_visuals)
            if title_shape is not None:
                fitted_title_size, _ = fit_title_font_size(
                    content["title"],
                    {
                        "left": int(title_shape.left),
                        "top": int(title_shape.top),
                        "width": int(title_shape.width),
                        "height": int(title_shape.height),
                    },
                    title_size,
                )
            else:
                fitted_title_size = title_size
            replace_shape_text(
                title_shape,
                content["title"],
                font_name=title_font,
                font_size=fitted_title_size,
                color=title_color,
                bold=True,
            )
            replace_shape_text(
                content_shape,
                "\n".join(points),
                bullet=content.get("content_mode") != "disclaimer",
                font_name=fonts["body_style"],
                font_size=body_size,
                color=body_color,
                bold=False,
            )
            align_bullet_text_with_title(content_shape, title_shape)
            replace_shape_text(
                source_shape,
                template_source_text(content) if allowed_source else "",
            )
            replace_shape_text(page_number_shape, str(plan["slide_number"]))
            visual_mappings = content.get("visual_mappings", [])
            if content_shape is not None and visual_mappings:
                required_body_height, estimated_lines = estimate_text_height(
                    "\n".join(points),
                    int(content_shape.width),
                    body_size,
                    bullet=True,
                )
                body_height = min(
                    1450000,
                    max(360000, required_body_height),
                )
                content_shape.height = Emu(body_height)
                visual_top = max(
                    visual_box["top"],
                    int(content_shape.top + content_shape.height) + 90000,
                )
                lower_boundaries = [
                    int(shape.top)
                    for shape in (source_shape, page_number_shape)
                    if shape is not None
                ]
                visual_bottom = (
                    min(lower_boundaries) - 90000
                    if lower_boundaries
                    else int(prs.slide_height) - 350000
                )
                visual_box["top"] = visual_top
                visual_box["height"] = max(900000, visual_bottom - visual_top)
                # 新增：视觉区域自动铺满页面
                visual_box["left"] = int(prs.slide_width * 0.05)
                visual_box["width"] = int(prs.slide_width * 0.90)

            single_visual_left = bullet_left_edge(content_shape)
            single_visual_box = (
                {
                    "left": single_visual_left,
                    "width": int(prs.slide_width) - 2 * single_visual_left,
                }
                if single_visual_left is not None
                else None
            )
            rendered_visual_count = add_source_visuals(
                slide,
                visual_mappings,
                document,
                visual_box,
                table_styles,
                Path(
                    data.get(
                        "asset_root",
                        Path(__file__).resolve().parent.parent,
                    )
                ),
                single_visual_box,
            )

        if layout_debug is not None:
            tracked_shapes: list[tuple[str, Any]] = [
                (name, shape)
                for name, shape in (
                    ("title", title_shape),
                    ("content", content_shape),
                    ("source", source_shape),
                    ("page_number", page_number_shape),
                )
                if shape is not None
            ]
            tracked_shapes.extend(
                (shape.name, shape)
                for shape in slide.shapes
                if shape.name.startswith("source_docx_")
                or shape.name == "source_table"
            )
            overlap_pairs: list[list[str]] = []
            for left_index, (left_name, left_shape) in enumerate(tracked_shapes):
                left_box = {
                    "left": int(left_shape.left),
                    "top": int(left_shape.top),
                    "width": int(left_shape.width),
                    "height": int(left_shape.height),
                }
                for right_name, right_shape in tracked_shapes[left_index + 1 :]:
                    right_box = {
                        "left": int(right_shape.left),
                        "top": int(right_shape.top),
                        "width": int(right_shape.width),
                        "height": int(right_shape.height),
                    }
                    if boxes_overlap(left_box, right_box):
                        overlap_pairs.append([left_name, right_name])
            layout_debug.append(
                {
                    "slide_number": plan["slide_number"],
                    "page_type": page_type,
                    "template_slide_number": mapped["slide_number"],
                    "creation_mode": "template_example_clone",
                    "rendered_visual_count": rendered_visual_count,
                    "actual_layout": {
                        name: {
                            "left": int(shape.left),
                            "top": int(shape.top),
                            "width": int(shape.width),
                            "height": int(shape.height),
                        }
                        for name, shape in tracked_shapes
                    },
                    "overlaps": overlap_pairs,
                    "placeholder_prompt_count": sum(
                        "单击" in (shape.text or "")
                        for shape in iter_ppt_shapes(slide.shapes)
                        if getattr(shape, "has_text_frame", False)
                    ),
                }
            )

    remove_initial_slides(prs, template_slide_count)
    
    return prs


def main() -> None:
    project_root = Path(__file__).resolve().parent.parent
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=project_root / "output",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=project_root / "output" / "generated_test.pptx",
    )
    parser.add_argument(
        "--layout-debug",
        type=Path,
        default=project_root / "output" / "layout_debug.json",
    )
    args = parser.parse_args()

    data = load_inputs(args.output_dir)
    layout_debug: list[dict[str, Any]] = []
    presentation = build_presentation(data, layout_debug)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(args.output)
    args.layout_debug.parent.mkdir(parents=True, exist_ok=True)
    args.layout_debug.write_text(
        json.dumps(
            {
                "slide_count": len(layout_debug),
                "overlap_count": sum(
                    len(item.get("overlaps", [])) for item in layout_debug
                ),
                "slides": layout_debug,
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    print(f"Wrote {args.output} ({len(presentation.slides)} slides)")
    print(f"Wrote {args.layout_debug}")


if __name__ == "__main__":
    main()
