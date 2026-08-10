from __future__ import annotations

import json
import tempfile
import unittest
import zipfile
from pathlib import Path

import yaml
from docx import Document as WordDocument
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.util import Inches

from scripts.generate_ppt import (
    FIXED_DISCLAIMER_PARAGRAPHS,
    build_presentation,
    dynamic_layout,
    fit_title_font_size,
    normalize_chart_axis_ids,
    order_generation_plans,
    resolve_template_path,
    replace_risk_content,
    replace_shape_text,
    replace_summary_text,
    select_template_mapping,
    table_cell_style,
    title_style_for_page,
)
from scripts.parse_templates import (
    heading_classification,
    paragraph_run_metadata,
    parse_docx,
    semantic_heading_level,
)
from scripts.plan_presentation_content import (
    build_plan as build_content_plan,
    compact_summary_paragraphs,
    display_source,
    extract_document_conclusions,
    map_visuals_to_preceding_summaries,
    normalize_cover_title,
    parse_sections as parse_content_sections,
    split_section_title,
    visual_points_for_slide,
)
from scripts.plan_slides import build_plan
from test_generate import count_layout_overlaps


PROJECT_ROOT = Path(__file__).resolve().parent.parent
FIXTURE = PROJECT_ROOT / "test" / "input" / "planner_heading_cases.docx"


class PlannerHeadingTests(unittest.TestCase):
    def test_chart_axis_ids_are_normalized_for_powerpoint(self) -> None:
        chart_xml = parse_xml(
            b"""
            <c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart">
              <c:chart>
                <c:axId val="-206191696"/>
                <c:crossAx val="-206191696"/>
              </c:chart>
            </c:chartSpace>
            """
        )
        normalize_chart_axis_ids(chart_xml)
        values = [
            int(element.get("val"))
            for element in chart_xml.iter()
            if element.tag.endswith(("axId", "crossAx"))
        ]
        self.assertEqual(values, [4088775600, 4088775600])
        normalize_chart_axis_ids(chart_xml, chart_namespace_id=18)
        remapped_values = [
            int(element.get("val"))
            for element in chart_xml.iter()
            if element.tag.endswith(("axId", "crossAx"))
        ]
        self.assertEqual(remapped_values, [2000000576, 2000000576])

    def test_second_real_report_metadata_and_numbered_headings(self) -> None:
        report = (
            PROJECT_ROOT
            / "test"
            / "input"
            / "【广发策略】位置决定叙事，趋势锚定景气——2026港股&海外大类资产中期策略V2.docx"
        )
        if not report.exists():
            self.skipTest("港股报告测试文件不存在")
        with tempfile.TemporaryDirectory() as temporary_directory:
            parsed = parse_docx(report, Path(temporary_directory))

        self.assertEqual(
            parsed["title"],
            "位置决定叙事，趋势锚定景气——2026港股&海外大类资产中期策略",
        )
        self.assertEqual(parsed["report_metadata"]["report_date"], "2026/07/12")
        self.assertEqual(
            parsed["report_metadata"]["authors"],
            ["刘晨明", "余可骋", "陈振威"],
        )
        self.assertEqual(len(parsed["report_metadata"]["summary_paragraphs"]), 9)
        self.assertEqual(
            [heading["text"] for heading in parsed["heading_1"]],
            [
                "一、全球宏观环境：滞胀下的鹰派转向",
                "二、港股流动性压力下半年能否缓解？",
                "三、大类资产走势复盘与研判",
                "四、大类资产配置总结",
                "五、风险提示",
            ],
        )

    def test_risk_page_uses_fixed_disclaimer_paragraphs(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(11), Inches(5)
        )
        replace_risk_content(shape, ["测试风险"], "输入文档中的旧免责声明")
        paragraphs = shape.text_frame.paragraphs
        self.assertEqual(
            [paragraph.text for paragraph in paragraphs[4:8]],
            list(FIXED_DISCLAIMER_PARAGRAPHS),
        )
        self.assertNotIn("旧免责声明", shape.text)

    def test_real_report_metadata_and_chart_inventory(self) -> None:
        report = (
            PROJECT_ROOT
            / "test"
            / "input"
            / "【广发策略】从杠杆繁荣到筹码松动：韩国杠杆去化走到哪一步？V3(1).docx"
        )
        if not report.exists():
            self.skipTest("韩国报告测试文件不存在")
        with tempfile.TemporaryDirectory() as temporary_directory:
            parsed = parse_docx(report, Path(temporary_directory))
            extracted_root = Path(temporary_directory).parent
            for chart in parsed["charts"]:
                self.assertEqual(chart["object_type"], "native_chart")
                self.assertEqual(chart["transfer_mode"], "native_editable")
                self.assertGreater(chart["width_emu"], 0)
                self.assertGreater(chart["height_emu"], 0)
                self.assertTrue(
                    (extracted_root / chart["workbook_path"]).is_file()
                )

        self.assertEqual(
            parsed["title"],
            "从杠杆繁荣到筹码松动：韩国杠杆去化走到哪一步？",
        )
        self.assertEqual(parsed["report_metadata"]["report_date"], "2026/07/20")
        self.assertEqual(parsed["report_metadata"]["authors"], ["刘晨明", "陈振威"])
        self.assertEqual(len(parsed["charts"]), 22)
        summary_paragraphs = parsed["report_metadata"]["summary_paragraphs"]
        self.assertEqual(len(summary_paragraphs), 6)
        self.assertTrue(summary_paragraphs[0]["text"].startswith("KOSPI本轮调整"))
        self.assertTrue(summary_paragraphs[-1]["text"].startswith("最后，我们构建了"))
        self.assertFalse(
            any("风险提示" in paragraph["text"] for paragraph in summary_paragraphs)
        )
        self.assertTrue(
            any(
                run["bold"]
                for paragraph in summary_paragraphs
                for run in paragraph["runs"]
            )
        )

    def test_generated_titles_do_not_add_continuation_suffix(self) -> None:
        content_path = PROJECT_ROOT / "output" / "presentation_content.json"
        if not content_path.exists():
            self.skipTest("presentation_content.json 尚未生成")
        content = json.loads(content_path.read_text(encoding="utf-8"))
        self.assertTrue(
            all("（续）" not in slide["title"] for slide in content["slides"])
        )

    def test_summary_copies_word_marker_paragraphs_with_emphasis(self) -> None:
        report = (
            PROJECT_ROOT
            / "test"
            / "input"
            / "【广发策略】从杠杆繁荣到筹码松动：韩国杠杆去化走到哪一步？V3(1).docx"
        )
        if not report.exists():
            self.skipTest("韩国报告测试文件不存在")
        with tempfile.TemporaryDirectory() as temporary_directory:
            document = parse_docx(report, Path(temporary_directory))
        presentation_plan = build_plan(document)
        content = build_content_plan(document, presentation_plan["slides"])
        page_types = {
            slide["slide_number"]: slide["page_type"]
            for slide in presentation_plan["slides"]
        }
        summary = next(
            slide
            for slide in content["slides"]
            if page_types[slide["slide_number"]] == "summary"
        )
        expected = compact_summary_paragraphs(
            document["report_metadata"]["summary_paragraphs"],
            maximum=5,
        )
        self.assertEqual(summary["rich_text"], expected)
        self.assertEqual(
            summary["key_points"],
            [paragraph["text"] for paragraph in expected],
        )
        self.assertFalse(any("风险提示" in text for text in summary["key_points"]))
        self.assertTrue(
            any(
                run["bold"]
                for paragraph in summary["rich_text"]
                for run in paragraph["runs"]
            )
        )

    def test_summary_renderer_keeps_blank_paragraphs_between_points(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(10), Inches(5)
        )
        shape.text_frame.text = "模板要点一\n\n模板要点二\n\n模板要点三"
        replace_summary_text(shape, ["甲", "乙", "丙"])
        paragraphs = shape.text_frame.paragraphs
        self.assertEqual([paragraphs[index].text for index in (0, 2, 4)], ["甲", "乙", "丙"])
        self.assertEqual([paragraphs[index].text for index in (1, 3)], ["", ""])
        self.assertEqual(paragraphs[0].line_spacing, 1.5)
        self.assertEqual(paragraphs[2].line_spacing, 1.5)
        self.assertEqual(paragraphs[4].line_spacing, 1.5)
        self.assertEqual(paragraphs[1].line_spacing, 1.5)

    def test_docx_run_metadata_detects_source_han_medium(self) -> None:
        document = WordDocument()
        paragraph = document.add_paragraph()
        medium_run = paragraph.add_run("图 1：Medium结论")
        medium_run._r.get_or_add_rPr().get_or_add_rFonts().set(
            qn("w:eastAsia"), "思源黑体 CN Medium"
        )
        regular_run = paragraph.add_run("普通补充")
        regular_run._r.get_or_add_rPr().get_or_add_rFonts().set(
            qn("w:eastAsia"), "思源黑体 CN Regular"
        )

        runs, _ = paragraph_run_metadata(paragraph)

        self.assertTrue(runs[0]["source_han_medium"])
        self.assertFalse(runs[1]["source_han_medium"])
        self.assertIn("思源黑体 CN Medium", runs[0]["font_names"])

    def test_body_renderer_copies_square_bullet_to_every_point(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(10), Inches(3)
        )
        paragraph = shape.text_frame.paragraphs[0]
        paragraph._p.get_or_add_pPr().append(
            parse_xml(
                b'<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="\xe2\x96\xaa"/>'
            )
        )
        paragraph.text = "模板正文"
        replace_shape_text(shape, "第一条\n第二条", bullet=True)
        for rendered in shape.text_frame.paragraphs:
            bullet_nodes = rendered._p.xpath("./a:pPr/a:buChar")
            self.assertEqual(len(bullet_nodes), 1)
            self.assertEqual(bullet_nodes[0].get("char"), "▪")

    def test_planner_omits_visual_free_body_pages(self) -> None:
        document = {
            "file": "sample.docx",
            "title": "示例报告",
            "images": [],
            "tables": [],
            "content_order": [
                {"type": "heading_1", "text": "一、正文", "order_index": 1},
                {"type": "heading_2", "text": "（一）只有文字", "order_index": 2},
                {"type": "body", "text": "这一小节没有任何视觉素材。", "order_index": 3},
                {"type": "heading_2", "text": "（二）包含图表", "order_index": 4},
                {"type": "body", "text": "图表对应的结论。", "order_index": 5},
                {"type": "chart", "index": 1, "chart_indexes": [1], "order_index": 6},
            ],
        }
        plan = build_plan(document)
        analytical = [
            slide
            for slide in plan["slides"]
            if slide["page_type"]
            not in {"cover", "summary", "section", "risk", "thanks"}
        ]
        self.assertEqual(len(analytical), 1)
        self.assertIn(6, analytical[0]["source_order_indexes"])
        self.assertNotIn(
            "（一）只有文字",
            [slide["source_heading"] for slide in analytical],
        )

    def test_visual_mapping_stops_at_previous_visual_boundary(self) -> None:
        items = [
            {
                "type": "body",
                "index": 1,
                "order_index": 1,
                "text": "普通说明。",
                "bold_sentences": ["加粗核心判断"],
                "runs": [
                    {
                        "text": "图 1：Medium核心判断",
                        "source_han_medium": True,
                    }
                ],
            },
            {
                "type": "table",
                "index": 1,
                "order_index": 2,
                "image_indexes": [11],
            },
            {
                "type": "body",
                "index": 2,
                "order_index": 3,
                "text": "第二个视觉对象前最近的解释句。后续细节。",
                "bold_sentences": [],
                "runs": [
                    {
                        "text": "图 2：第二个Medium结论",
                        "source_han_medium": True,
                    },
                    {
                        "text": "不应提取的普通文字",
                        "source_han_medium": False,
                    },
                ],
            },
            {
                "type": "image",
                "index": 12,
                "order_index": 4,
                "parent_type": "paragraph",
                "parent_index": 3,
            },
        ]

        mappings = map_visuals_to_preceding_summaries(items)

        self.assertEqual(len(mappings), 2)
        self.assertEqual(mappings[0]["paragraph_text"], "Medium核心判断。")
        self.assertEqual(mappings[0]["mapping_basis"], "source_han_medium")
        self.assertEqual(mappings[0]["table_id"], 1)
        self.assertEqual(mappings[0]["embedded_image_ids"], [11])
        self.assertEqual(mappings[1]["paragraph_text"], "第二个Medium结论。")
        self.assertEqual(mappings[1]["mapping_basis"], "source_han_medium")
        self.assertEqual(mappings[1]["image_id"], 12)

    def test_visual_mapping_does_not_cross_subsection_boundary(self) -> None:
        mappings = map_visuals_to_preceding_summaries(
            [
                {
                    "type": "body",
                    "index": 1,
                    "order_index": 1,
                    "text": "上一小节的核心结论已经明显改善。",
                    "bold_sentences": [],
                    "section_title": "第一章",
                    "subsection_title": "小节一",
                },
                {
                    "type": "table",
                    "index": 1,
                    "order_index": 2,
                    "rows": [["图 1：当前指标仍在下行"]],
                    "section_title": "第一章",
                    "subsection_title": "小节二",
                },
            ]
        )

        self.assertEqual(mappings[0]["mapping_basis"], "visual_title")
        self.assertEqual(mappings[0]["paragraph_index"], None)

    def test_page_level_visual_copy_deduplicates_and_adds_evidence(self) -> None:
        mappings = [
            {
                "paragraph_text": "共同结论。",
                "conclusion_text": "共同结论。",
                "evidence_text": "支持证据。",
                "evidence_score": {"total": 72.0},
            },
            {
                "paragraph_text": "共同结论。",
                "conclusion_text": "共同结论。",
                "evidence_text": "次要证据。",
                "evidence_score": {"total": 61.0},
            },
        ]

        self.assertEqual(
            visual_points_for_slide(mappings),
            ["共同结论。", "支持证据。"],
        )

    def test_multi_chart_container_keeps_one_title_per_chart(self) -> None:
        from scripts.plan_presentation_content import expand_visual_mappings

        mapped = map_visuals_to_preceding_summaries(
            [
                {
                    "type": "table",
                    "index": 1,
                    "order_index": 1,
                    "chart_indexes": [7, 8],
                    "rows": [["图 1：外资持仓下降", "图 2：散户承接增加"]],
                }
            ]
        )

        expanded = expand_visual_mappings(mapped)
        self.assertEqual(
            [mapping["conclusion_text"] for mapping in expanded],
            ["外资持仓下降。", "散户承接增加。"],
        )

    def test_visual_mapping_uses_container_text_when_no_paragraph_precedes(self) -> None:
        mappings = map_visuals_to_preceding_summaries(
            [
                {
                    "type": "table",
                    "index": 3,
                    "order_index": 1,
                    "rows": [
                        ["图 1：指数回撤"],
                        ["盈利预期转弱使拥挤交易快速松动。"],
                    ],
                }
            ]
        )

        self.assertEqual(len(mappings), 1)
        self.assertEqual(mappings[0]["paragraph_text"], "指数回撤。")
        self.assertEqual(mappings[0]["mapping_basis"], "visual_title")
        self.assertIn("score", mappings[0])
        self.assertTrue(mappings[0]["candidate_scores"])

    def test_visual_mapping_rejects_short_or_punctuation_only_bold_text(self) -> None:
        mappings = map_visuals_to_preceding_summaries(
            [
                {
                    "type": "body",
                    "index": 1,
                    "order_index": 1,
                    "text": "信用融资余额仍处于历史较高水平。",
                    "bold_sentences": ["。", "短句"],
                },
                {
                    "type": "image",
                    "index": 1,
                    "order_index": 2,
                    "parent_type": "paragraph",
                    "parent_index": 2,
                },
            ]
        )

        self.assertEqual(len(mappings), 1)
        self.assertEqual(mappings[0]["paragraph_text"], "")
        self.assertEqual(mappings[0]["mapping_basis"], "missing_candidate")

    def test_docx_content_order_contains_visual_parent_relationships(self) -> None:
        input_files = sorted((PROJECT_ROOT / "test" / "input").glob("*.docx"))
        real_reports = [
            path
            for path in input_files
            if path.name != "planner_heading_cases.docx"
            and not path.name.startswith("~$")
            and zipfile.is_zipfile(path)
        ]
        if not real_reports:
            self.skipTest("真实 DOCX 尚未放入 test/input")

        with tempfile.TemporaryDirectory() as temporary_directory:
            parsed = parse_docx(real_reports[0], Path(temporary_directory))

        content = parsed["content_order"]
        self.assertEqual(
            [item["order_index"] for item in content],
            list(range(1, len(content) + 1)),
        )
        images = [item for item in content if item["type"] == "image"]
        self.assertTrue(images)
        for image in images:
            self.assertIn(image["parent_type"], {"paragraph", "table"})
            self.assertIsInstance(image["parent_index"], int)
            parents = [
                item
                for item in content
                if item["type"] == image["parent_type"]
                and item["index"] == image["parent_index"]
            ]
            # An inline image can be the only content of an otherwise empty
            # paragraph; in that case the image entry itself preserves its
            # document position and parent paragraph index.
            if parents:
                self.assertLess(parents[-1]["order_index"], image["order_index"])

    def test_content_summary_source_section_and_risk_rules(self) -> None:
        doc_payload = json.loads(
            (PROJECT_ROOT / "output" / "doc_structure.json").read_text(
                encoding="utf-8"
            )
        )
        document = doc_payload["files"][0]
        presentation_plan = build_plan(document)
        content = build_content_plan(document, presentation_plan["slides"])
        planned_by_number = {
            slide["slide_number"]: slide for slide in presentation_plan["slides"]
        }

        summary = next(
            slide
            for slide in content["slides"]
            if planned_by_number[slide["slide_number"]]["page_type"] == "summary"
        )
        expected_summary = compact_summary_paragraphs(
            document["report_metadata"]["summary_paragraphs"],
            maximum=5,
        )
        self.assertEqual(
            summary["key_points"],
            [paragraph["text"] for paragraph in expected_summary],
        )
        self.assertEqual(summary["rich_text"], expected_summary)
        self.assertTrue(
            all(
                forbidden not in point
                for point in summary["key_points"]
                for forbidden in ("数据来源", "资料来源", ".docx")
            )
        )
        self.assertIsNone(summary["display_source"])

        source_pages = [
            slide
            for slide in content["slides"]
            if planned_by_number[slide["slide_number"]]["page_type"]
            in {"chart_analysis", "chart_table", "table_summary", "matrix"}
        ]
        self.assertTrue(any(slide["display_source"] for slide in source_pages))
        self.assertTrue(
            all(
                ".docx" not in (slide["display_source"] or "")
                for slide in source_pages
            )
        )

        section_two = next(
            slide
            for slide in content["slides"]
            if slide["source_trace"]["section"]
            and slide["source_trace"]["section"].startswith("二、")
            and planned_by_number[slide["slide_number"]]["page_type"] == "section"
        )
        self.assertEqual(section_two["title"], "韩国股市谁在加杠杆、谁在承接？")
        self.assertEqual(section_two["subtitle"], "三类杠杆与资金对手盘全景")

        risk_slides = [
            slide
            for slide in content["slides"]
            if planned_by_number[slide["slide_number"]]["page_type"] == "risk"
        ]
        self.assertTrue(risk_slides)
        self.assertTrue(
            all(
                slide["source_trace"]["section"].startswith("五、风险提示")
                for slide in risk_slides
            )
        )
        self.assertTrue(
            all(
                "四、如何识别二次冲击"
                not in slide["source_trace"]["section"]
                for slide in risk_slides
            )
        )

    def test_semantic_chapter_numbering_is_heading_one(self) -> None:
        for title in (
            "一、市场环境",
            "二、配置建议",
            "六、风险提示",
            "第七章 资产展望",
        ):
            with self.subTest(title=title):
                self.assertEqual(semantic_heading_level(title), 1)
        for title in ("1、宏观判断", "2．行业配置", "3. 策略建议"):
            with self.subTest(title=title):
                self.assertEqual(semantic_heading_level(title), 3)
        self.assertEqual(semantic_heading_level("（一）信用融资分析"), 2)
        self.assertEqual(
            semantic_heading_level("一、目录中的重复标题", None, "toc 1"),
            None,
        )
        self.assertEqual(
            heading_classification("二、真实章节", 2, "gfzq二级标题"),
            (1, "numbering"),
        )

    def test_semantic_numbered_body_items_create_sections(self) -> None:
        document = {
            "content_order": [
                {"type": "body", "text": "一、第一章", "style": "Normal"},
                {"type": "body", "index": 1, "text": "第一章正文。"},
                {"type": "heading_2", "text": "二、第二章", "style": "Heading 2"},
                {"type": "body", "index": 2, "text": "第二章正文。"},
            ]
        }
        sections = parse_content_sections(document)
        self.assertEqual(
            [section["title"] for section in sections],
            ["一、第一章", "二、第二章"],
        )

    def test_every_non_risk_heading_one_creates_section_page(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            document = parse_docx(FIXTURE, Path(temp_dir))
        plan = build_plan(document)
        heading_titles = [item["text"] for item in document["heading_1"]]
        section_pages = [
            slide for slide in plan["slides"] if slide["page_type"] == "section"
        ]

        self.assertEqual(len(heading_titles), 3)
        expected_section_titles = [
            title for title in heading_titles if "风险提示" not in title
        ]
        self.assertEqual(len(section_pages), len(expected_section_titles))
        self.assertEqual(
            [slide["source_heading"] for slide in section_pages],
            expected_section_titles,
        )
        self.assertTrue(
            any(
                slide["page_type"] == "section"
                and "无二级标题" in slide["source_heading"]
                for slide in plan["slides"]
            )
        )
        self.assertFalse(
            any(
                slide["page_type"] == "section"
                and "风险提示" in slide["source_heading"]
                for slide in plan["slides"]
            )
        )
        self.assertTrue(any(slide["page_type"] == "risk" for slide in plan["slides"]))
        self.assertEqual(plan["slides"][-1]["page_type"], "thanks")

    def test_cover_title_preserves_source_wording(self) -> None:
        self.assertEqual(
            normalize_cover_title(
                "从杠杆繁荣到筹码松动：韩国杠杆去化走到哪一步？"
            ),
            "从杠杆繁荣到筹码松动：韩国杠杆去化走到哪一步？",
        )

    def test_content_planner_accepts_all_section_pages(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            document = parse_docx(FIXTURE, Path(temp_dir))
        presentation_plan = build_plan(document)
        content_plan = build_content_plan(document, presentation_plan["slides"])
        section_slides = [
            slide
            for slide, planned in zip(
                content_plan["slides"],
                presentation_plan["slides"],
            )
            if planned["page_type"] == "section"
        ]

        non_risk_headings = [
            heading
            for heading in document["heading_1"]
            if "风险提示" not in heading["text"]
        ]
        self.assertEqual(len(section_slides), len(non_risk_headings))
        self.assertEqual(
            [slide["source_trace"]["section"] for slide in section_slides],
            [heading["text"] for heading in non_risk_headings],
        )

    def test_heading_level_one_is_always_a_section(self) -> None:
        sections = parse_content_sections(
            {
                "content_order": [
                    {
                        "type": "heading",
                        "heading_level": 1,
                        "text": "无二级标题的一级章节",
                    },
                    {"type": "body", "index": 1, "text": "章节正文。"},
                ]
            }
        )
        self.assertEqual(len(sections), 1)
        self.assertEqual(sections[0]["title"], "无二级标题的一级章节")


class GeneratorStyleTests(unittest.TestCase):
    def test_layout_debug_overlap_count_tolerates_missing_field(self) -> None:
        layout_debug = [
            {"slide_number": 1, "overlaps": [["title", "content"]]},
            {"slide_number": 2},
            {"slide_number": 3, "overlaps": []},
            {
                "slide_number": 4,
                "overlaps": [["content", "chart"], ["chart", "source"]],
            },
        ]
        self.assertEqual(count_layout_overlaps(layout_debug), 3)

    def test_template_clone_has_no_prompt_or_forbidden_page_text(self) -> None:
        output = PROJECT_ROOT / "output"
        plan_slides = [
            {"slide_number": 1, "page_type": "cover"},
            {"slide_number": 2, "page_type": "summary"},
            {"slide_number": 3, "page_type": "section"},
            {"slide_number": 4, "page_type": "chart_analysis"},
        ]
        content_slides = [
            {
                "slide_number": 1,
                "title": "测试封面",
                "key_points": ["COVER_BODY_MUST_NOT_APPEAR"],
                "source": {"document": "test.docx"},
            },
            {
                "slide_number": 2,
                "title": "核心结论",
                "key_points": ["第一条核心观点", "第二条核心观点"],
                "source": {"document": "test.docx"},
            },
            {
                "slide_number": 3,
                "title": "测试章节",
                "key_points": ["SECTION_BODY_MUST_NOT_APPEAR"],
                "source": {"document": "test.docx"},
            },
            {
                "slide_number": 4,
                "title": "正文分析",
                "key_points": ["正文观点"],
                "source": {"document": "test.docx", "section": "测试章节"},
            },
        ]
        style = yaml.safe_load(
            (output / "style_config.yaml").read_text(encoding="utf-8")
        )
        data = {
            "plan": {
                "planned_slide_count": len(plan_slides),
                "slides": plan_slides,
            },
            "content": {"slides": content_slides},
            "mapping": yaml.safe_load(
                (output / "slide_mapping.yaml").read_text(encoding="utf-8")
            ),
            "layout": json.loads(
                (output / "ppt_layout.json").read_text(encoding="utf-8")
            ),
            "style": style,
            "template_path": PROJECT_ROOT / "template" / style["source_template"],
        }
        debug: list[dict[str, object]] = []
        presentation = build_presentation(data, debug)

        slide_texts = [
            [
                shape.text
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
            ]
            for slide in presentation.slides
        ]
        all_text = "\n".join(text for texts in slide_texts for text in texts)
        self.assertNotIn("单击添加标题", all_text)
        self.assertNotIn("单击此处编辑母版文本样式", all_text)
        self.assertNotIn("COVER_BODY_MUST_NOT_APPEAR", "\n".join(slide_texts[0]))
        self.assertNotIn("数据来源", "\n".join(slide_texts[1]))
        self.assertNotIn("数据来源", "\n".join(slide_texts[2]))
        self.assertNotIn("SECTION_BODY_MUST_NOT_APPEAR", "\n".join(slide_texts[2]))
        self.assertIn("第一条核心观点", "\n".join(slide_texts[1]))
        self.assertTrue(
            all(item["creation_mode"] == "template_example_clone" for item in debug)
        )

    def test_section_mapping_is_always_template_slide_three(self) -> None:
        examples = {
            "section": [
                {"slide_number": 3, "page_type": "section"},
                {"slide_number": 7, "page_type": "section"},
                {"slide_number": 14, "page_type": "section"},
            ]
        }
        usage: dict[str, int] = {}
        selected = [
            select_template_mapping("section", examples, usage)
            for _ in range(5)
        ]
        self.assertEqual([item["slide_number"] for item in selected], [3] * 5)

    def test_thanks_mapping_uses_semantic_mapping_instead_of_fixed_number(self) -> None:
        examples = {
            "thanks": [
                {"slide_number": 1, "page_type": "thanks"},
            ]
        }
        selected = select_template_mapping("thanks", examples, {})
        self.assertEqual(selected["slide_number"], 1)

    def test_thanks_plan_is_always_generated_last(self) -> None:
        plans = [
            {"slide_number": 1, "page_type": "thanks"},
            {"slide_number": 30, "page_type": "risk"},
            {"slide_number": 2, "page_type": "cover"},
            {"slide_number": 10, "page_type": "chart_analysis"},
        ]
        ordered = order_generation_plans(plans)
        self.assertEqual(
            [(item["slide_number"], item["page_type"]) for item in ordered],
            [
                (2, "cover"),
                (10, "chart_analysis"),
                (30, "risk"),
                (1, "thanks"),
            ],
        )

    def test_template_path_resolves_to_original_pptx(self) -> None:
        source_name = "【广发策略】全球流动性走到哪里了？——央.pptx"
        resolved = resolve_template_path(
            {"style": {"source_template": source_name}}
        )
        self.assertEqual(resolved.name, source_name)
        self.assertTrue(resolved.is_file())

    def test_section_and_body_title_sizes_never_mix(self) -> None:
        fonts = {
            "cover_title_style": "Cover",
            "section_title_style": "Section",
            "title_style": "Body",
        }
        sizes = {
            "cover_title_pt": 38,
            "section_title_pt": 40,
            "body_title_pt": 32,
        }
        self.assertEqual(
            title_style_for_page("section", fonts, sizes),
            ("Section", 40.0),
        )
        self.assertEqual(
            title_style_for_page("chart_analysis", fonts, sizes),
            ("Body", 32.0),
        )

    def test_normal_title_font_shrinks_to_configured_box(self) -> None:
        box = {"left": 0, "top": 0, "width": 2_500_000, "height": 500_000}
        size, _ = fit_title_font_size(
            "这是一个需要自动缩小字号以避免超过标题框高度的普通页面长标题",
            box,
            32,
        )
        self.assertLess(size, 32)
        self.assertGreaterEqual(size, 1)

    def test_section_title_keeps_configured_size_and_expands_for_two_lines(self) -> None:
        title_seed = {
            "left": 600_000,
            "top": 500_000,
            "width": 3_000_000,
            "height": 500_000,
        }
        source = {
            "left": 500_000,
            "top": 6_200_000,
            "width": 5_000_000,
            "height": 250_000,
        }
        footer = {
            "left": 9_000_000,
            "top": 6_200_000,
            "width": 2_000_000,
            "height": 250_000,
        }
        regions, debug = dynamic_layout(
            page_type="section",
            title="这是一个用于验证章节标题支持两行显示的较长一级标题",
            body_text="",
            title_seed=title_seed,
            content_seed={
                "left": 600_000,
                "top": 1_300_000,
                "width": 10_000_000,
                "height": 800_000,
            },
            source_box=source,
            footer_box=footer,
            slide_width=12_192_000,
            title_size=40,
            body_size=14,
            has_chart=False,
            has_table=False,
            expand_title_height=True,
        )
        self.assertGreaterEqual(debug["title_estimated_lines"], 2)
        self.assertGreater(regions["title"]["height"], title_seed["height"])

    def test_table_regions_read_independent_styles(self) -> None:
        style_values = {
            "table": {
                "header": {
                    "font_name": "Header Font",
                    "font_size_pt": 10,
                    "text_color": "#FFFFFF",
                    "fill_color": "#353160",
                    "bold": True,
                    "alignment": "center",
                },
                "first_column": {
                    "font_name": "First Column Font",
                    "font_size_pt": 9,
                    "text_color": "#353160",
                    "fill_color": "#EEEEEE",
                    "bold": True,
                    "alignment": "left",
                },
                "body": {
                    "font_name": "Body Font",
                    "font_size_pt": 8,
                    "text_color": "#000000",
                    "fill_color": "#FFFFFF",
                    "bold": False,
                    "alignment": "right",
                },
            }
        }
        common = {
            "fallback_font": "Fallback",
            "fallback_size": 7,
            "fallback_text_color": RGBColor(0, 0, 0),
            "fallback_fill_color": None,
            "fallback_bold": False,
            "fallback_alignment": PP_ALIGN.CENTER,
        }
        header = table_cell_style(style_values, "header", **common)
        first_column = table_cell_style(style_values, "first_column", **common)
        body = table_cell_style(style_values, "body", **common)

        self.assertEqual(header["font_name"], "Header Font")
        self.assertEqual(first_column["font_name"], "First Column Font")
        self.assertEqual(body["font_name"], "Body Font")
        self.assertNotEqual(header["fill_color"], first_column["fill_color"])
        self.assertEqual(first_column["alignment"], PP_ALIGN.LEFT)
        self.assertEqual(body["alignment"], PP_ALIGN.RIGHT)


if __name__ == "__main__":
    unittest.main()
